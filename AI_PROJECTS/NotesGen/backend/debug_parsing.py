#!/usr/bin/env python3
"""Debug script to analyze what the enhanced parser is detecting."""

import sys
import os
sys.path.append('.')

from app.utils.enhanced_speaker_notes_parser import EnhancedSpeakerNotesParser
import xml.etree.ElementTree as ET

def debug_parsing():
    # Use the uploaded PPT file
    ppt_file = "uploads/ILT-TF-200-DEA-10-EN_M02.pptx"
    
    if not os.path.exists(ppt_file):
        print(f"❌ PPT file not found: {ppt_file}")
        return
    
    print(f"🔍 Debugging parsing for: {ppt_file}")
    
    parser = EnhancedSpeakerNotesParser()
    
    try:
        # Get raw XML directly
        from pptx import Presentation
        prs = Presentation(ppt_file)
        slide = prs.slides[0]
        
        if slide.has_notes_slide:
            notes_slide = slide.notes_slide
            xml_content = notes_slide._element.xml
            
            # First, let's see the raw text content
            notes_text = notes_slide.notes_text_frame.text
            
            print("\n" + "="*50)
            print("RAW NOTES TEXT:")
            print("="*50)
            print(repr(notes_text))  # Show exact string with escape chars
            
            # Parse XML to see paragraph structure
            root = ET.fromstring(xml_content)
            paragraphs = root.findall('.//a:p', parser.NAMESPACES)
            
            print(f"\n" + "="*50)
            print(f"FOUND {len(paragraphs)} XML PARAGRAPHS:")
            print("="*50)
            
            for i, para in enumerate(paragraphs):
                # Extract text from this paragraph
                text_elements = para.findall('.//a:t', parser.NAMESPACES)
                para_text = ''.join([elem.text or '' for elem in text_elements])
                
                print(f"\nParagraph {i+1}:")
                print(f"  Raw text: {repr(para_text)}")  # Show with escape chars
                print(f"  Length: {len(para_text)}")
                
                # Check if it would be detected as section header
                is_header, section_type = parser._identify_section_header(para_text)
                print(f"  Is section header: {is_header}, Type: {section_type}")
                
                # Show line breaks and structure
                lines = para_text.split('\n')
                if len(lines) > 1:
                    print(f"  Contains {len(lines)} lines:")
                    for j, line in enumerate(lines):
                        print(f"    Line {j+1}: {repr(line)}")
            
            # Test section detection on individual lines
            print(f"\n" + "="*50)
            print("TESTING SECTION DETECTION ON LINES:")
            print("="*50)
            
            all_lines = notes_text.split('\n')
            for i, line in enumerate(all_lines):
                line = line.strip()
                if line:
                    is_header, section_type = parser._identify_section_header(line)
                    if is_header:
                        print(f"✅ Line {i+1} IS HEADER: {repr(line)} -> {section_type}")
                    else:
                        print(f"   Line {i+1}: {repr(line)}")
            
        else:
            print("❌ No notes slide found")
            
    except Exception as e:
        print(f"❌ Error: {e}")
        import traceback
        traceback.print_exc()

if __name__ == '__main__':
    debug_parsing() 