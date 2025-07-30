#!/usr/bin/env python3
"""Debug script to examine raw text character by character."""

import sys
import os
sys.path.append('.')

def debug_raw():
    # Use the uploaded PPT file
    ppt_file = "uploads/ILT-TF-200-DEA-10-EN_M02.pptx"
    
    if not os.path.exists(ppt_file):
        print(f"❌ PPT file not found: {ppt_file}")
        return
    
    print(f"🔍 Debugging raw text for: {ppt_file}")
    
    try:
        # Get raw text from the presentation
        from pptx import Presentation
        prs = Presentation(ppt_file)
        slide = prs.slides[0]
        
        if slide.has_notes_slide:
            notes_slide = slide.notes_slide
            
            # Get raw text using python-pptx
            raw_text = notes_slide.notes_text_frame.text
            
            print(f"\n" + "="*50)
            print("RAW TEXT CHARACTER ANALYSIS:")
            print("="*50)
            
            print(f"Text length: {len(raw_text)}")
            print(f"Text repr: {repr(raw_text[:200])}...")
            
            # Find all unique characters
            unique_chars = set(raw_text)
            special_chars = [c for c in unique_chars if ord(c) < 32 or ord(c) > 126]
            
            print(f"\nSpecial characters found:")
            for char in sorted(special_chars, key=ord):
                print(f"  {repr(char)} (ord {ord(char)})")
            
            # Split by different possible line separators
            splits = {
                '\\n': raw_text.split('\n'),
                '\\x0b': raw_text.split('\x0b'), 
                '\\r': raw_text.split('\r'),
                '\\r\\n': raw_text.split('\r\n'),
            }
            
            print(f"\nSplit results:")
            for sep, parts in splits.items():
                print(f"  Split by {sep}: {len(parts)} parts")
                if len(parts) > 1:
                    for i, part in enumerate(parts[:5]):  # Show first 5 parts
                        print(f"    Part {i+1}: {repr(part[:50])}")
                    if len(parts) > 5:
                        print(f"    ... and {len(parts) - 5} more parts")
            
            # Look for section headers in the raw text
            print(f"\nSearching for section patterns:")
            instructor_pos = raw_text.find('|INSTRUCTOR NOTES:')
            student_pos = raw_text.find('|STUDENT NOTES:')
            
            print(f"  |INSTRUCTOR NOTES: found at position {instructor_pos}")
            print(f"  |STUDENT NOTES: found at position {student_pos}")
            
            if instructor_pos >= 0 and student_pos >= 0:
                instructor_section = raw_text[instructor_pos:student_pos]
                student_section = raw_text[student_pos:]
                
                print(f"\nInstructor section ({len(instructor_section)} chars):")
                print(f"  {repr(instructor_section)}")
                
                print(f"\nStudent section ({len(student_section)} chars):")
                print(f"  {repr(student_section[:200])}...")
            
        else:
            print("❌ No notes slide found")
            
    except Exception as e:
        print(f"❌ Error: {e}")
        import traceback
        traceback.print_exc()

if __name__ == '__main__':
    debug_raw() 