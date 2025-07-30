# ascendAWSClientManager.py
# Forge Ascend v4.1
# Updated by Tom Stern, 16 DEC 2024
#
#   based on Ascend 1 -- first version January 22 2024 -- by Tom Stern
#
#   1. Handles credentials
#   2. Creates a dictionary of clients to AWS Services and intializes them
#   3. Initializes clients and reinitializes clients
#
#   self.aws is the dictionary containing clients. The calling program must define:  self.aws = {}
#
#     Methods
#
#
import sys
import os
import json
import boto3
from botocore.exceptions import ClientError

from PyQt5.QtCore import Qt, QItemSelectionModel
from PyQt5.QtGui import QFont, QStandardItemModel, QStandardItem, QIcon, QTextCharFormat
from PyQt5.QtWidgets import (QApplication, QWidget, QHBoxLayout, QVBoxLayout, QLabel,
                             QTextEdit, QPushButton, QFrame, QTreeView, QMainWindow,
                             QRadioButton, QGridLayout, QGroupBox, QFormLayout,QInputDialog,QListWidget,
                             QFileDialog, QDialog, QMessageBox, QLineEdit, QStyle,QAbstractItemView)
import io
from docx import Document
import fitz  # Import PyMuPDF
import openpyxl
import csv
import json
import markdown
from pptx import Presentation
from bs4 import BeautifulSoup
import tempfile

## AWS Credentials
#   
#     ascendAWSClientManager
#
#   Ascend credentials follows the philosophy of putting the control in the hands of the user
#   Everything is done by button press rather than automatically, so the user can select which
#   credentials to use.
#   
#   This code is designed to accept JSON token-based credentials from Isengard
#   Standard credentials from a separate customer AWS account
#   Or from OS environment variables
#   
#   This is how you set the credentials in environment variables 
#   
#   Windows Command Line
#   set AWS_ACCESS_KEY_ID=your_access_key_id
#   set AWS_SECRET_ACCESS_KEY=your_secret_access_key
#   set AWS_SESSION_TOKEN=your_session_token
#   set AWS_DEFAULT_REGION=your_region
#   
#   Windows Powershell
#   $env:AWS_ACCESS_KEY_ID="your_access_key_id"
#   $env:AWS_SECRET_ACCESS_KEY="your_secret_access_key"
#   $env:AWS_SESSION_TOKEN="your_session_token"
#   $env:AWS_DEFAULT_REGION="your_region"
#   
#   MacOS and Linux
#   export AWS_ACCESS_KEY_ID=your_access_key_id
#   export AWS_SECRET_ACCESS_KEY=your_secret_access_key
#   export AWS_SESSION_TOKEN=your_session_token
#   export AWS_DEFAULT_REGION=your_region


# ------------ Provide a Control Panel for the user to control AWS credentials ----------- 
#
class CredentialsDialog(QDialog): ## --- Optional Session Token, Line edit entry
    def __init__(self, credentials, default_region, view_environment, parent=None):
        super().__init__(parent)
        self.credentials = credentials
        self.default_region = default_region
        self.view_environment = view_environment
        self.initUI()

    def initUI(self):
        self.setWindowTitle("Credentials, Session, and Client Controls")
                # ------------------[ Coloring Kit ]---------------------
        self.setStyleSheet("background-color: #E6E6E6; color: black;")
        self.buttonStyle_9 = """
        QPushButton {
            background-color: #C4E0EF;
            color: #000000;
            font-family: Arial; 
            font-size: 12px;    
            font-weight: normal;  
            font-style: normal; 
            border: 1px solid #000000;
            border-radius: 2px;
            }
            QPushButton:hover { background-color: #5b5b5b; color: #FFFFFF;}
            QPushButton:pressed { background-color: #FF0000; color: #000000; }
        """
        self.lineStyle_9 = """
        QLineEdit{
            background-color: #FFFFFF;
            color: #000000;
            font-family: Arial; 
            font-size: 12px;    
            font-weight: normal;  
            font-style: normal; 
            border: 1px solid #000000; }
        """
        self.textStyle_9 = """
        QTextEdit{
            background-color: #FFFFFF;
            color: #000000;
            font-family: Arial; 
            font-size: 12px;    
            font-weight: normal;  
            font-style: normal; 
            border: 1px solid #000000; }
        """
        #  self.Display.setStyleSheet(self.lineStyle_9)
        #  self.Button.setStyleSheet(self.buttonStyle_9)
        #  ------------------[ Coloring Kit ]---------------------


        credlayout = QVBoxLayout(self)
        credlayout_titles = QHBoxLayout()
        topH_layout = QHBoxLayout()
        layout_text = QVBoxLayout()
        layout_text_r1 = QHBoxLayout()
        layout_text_r2 = QHBoxLayout()
        layout_text_r3 = QHBoxLayout()
        layout_text_r4 = QHBoxLayout()
        layout_json = QVBoxLayout()
        layout_buttons = QHBoxLayout()
        layout_buttons_r2= QHBoxLayout()

        LW = 200    
               
        # Text box to retrieve JSON credentials with token
        self.json_edit = QTextEdit()
        self.json_edit.setStyleSheet(self.textStyle_9)
        self.json_edit.setFixedWidth(LW+100) 
        layout_json.addWidget(self.json_edit)

        # Retrieve these values from the environment, if they exist.
        self.env_aws_access_key_id = os.getenv('AWS_ACCESS_KEY_ID', '')
        self.env_aws_secret_access_key = os.getenv('AWS_SECRET_ACCESS_KEY', '')
        self.env_aws_session_token = os.getenv('AWS_SESSION_TOKEN', '')
        self.env_aws_region = os.getenv('AWS_DEFAULT_REGION', self.default_region)

        # Line edits for manual credential input
        self.accessKeyIDEdit = QLineEdit(self)
        self.accessKeyIDEdit.setFixedWidth(LW) 
        self.accessKeyIDEdit.setStyleSheet(self.lineStyle_9)
        self.secretAccessKeyEdit = QLineEdit(self)
        self.secretAccessKeyEdit.setFixedWidth(LW) 
        self.secretAccessKeyEdit.setStyleSheet(self.lineStyle_9)
        self.sessionToken = QLineEdit(self)
        self.sessionToken.setFixedWidth(LW) 
        self.sessionToken.setStyleSheet(self.lineStyle_9)
        self.regionEdit = QLineEdit(self)
        self.regionEdit.setFixedWidth(LW) 
        self.regionEdit.setStyleSheet(self.lineStyle_9)

        if self.env_aws_access_key_id and self.view_environment == 1:
            self.accessKeyIDEdit.setText(self.env_aws_access_key_id)   
        else:
            self.accessKeyIDEdit.setPlaceholderText("Enter Access Key ID") 
        if self.env_aws_secret_access_key and self.view_environment == 1:
            self.secretAccessKeyEdit.setText(self.env_aws_secret_access_key)   
        else:
            self.secretAccessKeyEdit.setPlaceholderText("Enter Secret Access Key") 
        if self.env_aws_session_token and self.view_environment == 1:
            self.sessionToken.setText(self.env_aws_session_token)   
        else:
            self.sessionToken.setPlaceholderText("Enter Session Token")               
        if self.env_aws_region and self.view_environment == 1:
            self.regionEdit.setText(self.env_aws_region)   
        else:
            self.regionEdit.setPlaceholderText("Enter AWS Region") 

        # Accept Button
        accept_button1 = QPushButton("Accept Creds")
        accept_button1.setStyleSheet(self.buttonStyle_9)
        accept_button1.clicked.connect(self.acceptCredentials)
        layout_buttons.addWidget(accept_button1)


        # Accept JSON Button
        accept_button2 = QPushButton("Accept JSON")
        accept_button2.setStyleSheet(self.buttonStyle_9)
        accept_button2.clicked.connect(self.acceptJSONCredentials) ##
        layout_buttons.addWidget(accept_button2)

        # Cancel
        accept_button3 = QPushButton("Cancel Without Changes")
        accept_button3.setStyleSheet(self.buttonStyle_9)
        accept_button3.clicked.connect(self.acceptNothing) ##
        layout_buttons_r2.addWidget(accept_button3)
    
        self.label1 = QLabel("Access Key ID:")
        layout_text_r1.addWidget(self.label1)
        layout_text_r1.addWidget(self.accessKeyIDEdit)
        self.label2 = QLabel("Secret Key:   ")
        layout_text_r2.addWidget(self.label2)
        layout_text_r2.addWidget(self.secretAccessKeyEdit)
        self.label3 = QLabel("Session Token:")
        layout_text_r3.addWidget(self.label3)
        layout_text_r3.addWidget(self.sessionToken)
        self.label4 = QLabel("Region:       ")
        layout_text_r4.addWidget(self.label4)
        layout_text_r4.addWidget(self.regionEdit)

        self.label5 = QLabel("Environment or TEXT Credentials:")
        self.label6 = QLabel("JSON Credentials:")
        credlayout_titles.addWidget(self.label5)
        credlayout_titles.addWidget(self.label6)

        credlayout.addLayout(credlayout_titles)
        layout_text.addLayout(layout_text_r1)
        layout_text.addLayout(layout_text_r2)
        layout_text.addLayout(layout_text_r3)
        layout_text.addLayout(layout_text_r4)
        topH_layout.addLayout(layout_text)
        topH_layout.addLayout(layout_json)
        credlayout.addLayout(topH_layout)
        credlayout.addLayout(layout_buttons)
        credlayout.addLayout(layout_buttons_r2)

    def acceptNothing(self):
        # Don't change the credentials, just close the dialog
        self.accept()

    def acceptCredentials(self):
        # Store the entered credentials and close the dialog
        ## Fill in the blanks with environment values
        try:
            self.accessKeyID = self.accessKeyIDEdit.text()
            if self.accessKeyID == "":
                self.accessKeyID = self.env_aws_access_key_id
            self.secretAccessKey = self.secretAccessKeyEdit.text()
            if self.secretAccessKey  == "":
                self.secretAccessKey  = self.env_aws_secret_access_key          
            if not self.accessKeyID or not self.secretAccessKey:
                raise ValueError("Both Access Key ID and Secret Access Key must be provided")
            self.sessionToken = self.sessionToken.text()
            if self.sessionToken  == "":
                self.sessionToken  = self.env_aws_session_token  
            self.regionName = self.regionEdit.text()
            if self.regionName  == "":
                self.regionName  = self.env_aws_region 
            if not self.regionName:
                self.regionName = self.default_region

            self.credentials['ACCESS'] = self.accessKeyID
            self.credentials['SECRET'] = self.secretAccessKey
            self.credentials['TOKEN'] = self.sessionToken
            self.credentials['REGION'] = self.regionName
            self.accept()
        except Exception as e:
            print("Error processing credentials:", e)
            QMessageBox.critical(self, "Error", "Issue processing credentials. Both Access Key ID and Secret Access Key must be provided.")

    def acceptJSONCredentials(self):
        # Parse and store the JSON credentials and close the dialog
        try:
            decode_credentials = json.loads(self.json_edit.toPlainText())
            self.accessKeyID = decode_credentials["credentials"]["accessKeyId"]
            self.secretAccessKey = decode_credentials["credentials"]["secretAccessKey"]
            self.sessionToken = decode_credentials["credentials"]["sessionToken"]
            # Pull region from Environment or from text box -- verify if Region is set in JSON creds in Isengard
            # self.regionName = decode_credentials["credentials"]["regionName"]
            self.regionName = self.regionEdit.text()
            if not self.regionName:
                self.regionName = self.default_region
            
            self.credentials['ACCESS'] = self.accessKeyID
            self.credentials['SECRET'] = self.secretAccessKey
            self.credentials['TOKEN'] = self.sessionToken
            self.credentials['REGION'] = self.regionName
            self.accept()
        except Exception as e:
            print("Error parsing credentials:", e)
            QMessageBox.critical(self, "Error", "Issue parsing credentials")

## ------------------ End of Custom Dialogs ----------------------



# Override the model to ensure that names are not editable in a QTreeView viewer
class CustomStandardItemModel(QStandardItemModel):
    def flags(self, index):
        default_flags = super(CustomStandardItemModel, self).flags(index)
        # Remove the Qt.ItemIsEditable flag to make items not editable
        return default_flags & ~Qt.ItemIsEditable
    

## ------------------ Class Definition --------------------------- [CLIENT MANAGER]
#
#
class ascendAWSClientManager:
    def __init__(self):
        #  self.clients = {}   This is the dictionary used globally to hold all AWS service clients Defined in main program
        self.session = None         # Class scoped session. Clients are created from the single session
        self.credentials = {}       # Class scoped credentials dictionary

        self.default_region = "us-east-1"  # <-- This is where you set the default region ## DEFAULT
        self.view_environment =  1         # <-- Set to ZERO to hide environment values from the user but use them if user did not enter values

    def load_credentials(self,clients):
        ## -- Open panel and get credentials from user or environment
        dialog = CredentialsDialog(self.credentials,self.default_region,self.view_environment)
        dialog.exec_()
        if 'ACCESS' in self.credentials or 'SECRET' in self.credentials:
                    self.initialize_aws_session()
                    self.add_clients(clients)


    def initialize_aws_session(self):
        aws_access_key_id = self.credentials['ACCESS']
        aws_secret_access_key = self.credentials['SECRET']
        aws_region_name = self.credentials['REGION']
        aws_session_token = self.credentials['TOKEN']

        if aws_session_token == "":
            # print("NO Token")
            self.session = boto3.Session(
                aws_access_key_id=aws_access_key_id,
                aws_secret_access_key=aws_secret_access_key,
                region_name=aws_region_name
            )
        else:
            # print("Token")
            self.session = boto3.Session(
                aws_access_key_id=aws_access_key_id,
                aws_secret_access_key=aws_secret_access_key,
                aws_session_token=aws_session_token,
                region_name=aws_region_name
            )

# clients['xyz'] is how all services can be initiated and passed around in the application within one ClientManager object.
#
#
    def add_clients(self,clients):
        clients['s3'] = self.session.client('s3') ## Example for future use
        clients['bedrock'] = self.session.client(service_name='bedrock')
        clients['bedrun'] = self.session.client(service_name='bedrock-runtime')
        clients['translate'] = self.session.client(service_name='translate')
        clients['polly'] = self.session.client(service_name='polly')
        # return 



### ==========================================================================  AWS Service MANAGERS

### ==========================================================================  AWS Service MANAGERS

# ------------ Provide a Control Panel for S3 Interaction ----------- 
#

## ------------------ Class Definition ---------------------------[ S3 MANAGER ]
#
#

    def ascendS3Manager(self, clients, edit_1=None, edit_2=None, edit_3=None):
        self.edit_1 = edit_1
        self.edit_2 = edit_2
        self.edit_3 = edit_3
        # Store clients in the instance if needed elsewhere
        self.clients = clients
        if 's3' not in self.clients:
            QMessageBox.warning(None, 'Warning', 'Credentials Issue. Could not use S3.')
            return
        # Create and setup the dialog
        self.dialog = QDialog()
        self.dialog.setWindowTitle("S3 Transporter")
        self.dialog.setGeometry(100, 100, 1200, 800)

        # ------------------[ Coloring Kit ]---------------------
        self.dialog.setStyleSheet("background-color: #E6E6E6; color: black;")
        self.buttonStyle_9 = """
        QPushButton {
            background-color: #C4E0EF;
            color: #000000;
            font-family: Arial; 
            font-size: 12px;    
            font-weight: normal;  
            font-style: normal; 
            border: 1px solid #000000;
            border-radius: 2px;
            }
            QPushButton:hover { background-color: #5b5b5b; color: #FFFFFF;}
            QPushButton:pressed { background-color: #FF0000; color: #000000; }
        """
        self.lineStyle_9 = """
        QLineEdit{
            background-color: #FFFFFF;
            color: #000000;
            font-family: Arial; 
            font-size: 12px;    
            font-weight: normal;  
            font-style: normal; 
            border: 1px solid #000000; }
        """
        self.textStyle_9 = """
        QTextEdit{
            background-color: #FFFFFF;
            color: #000000;
            font-family: Arial; 
            font-size: 12px;    
            font-weight: normal;  
            font-style: normal; 
            border: 1px solid #000000; }
        """
        #  self.Display.setStyleSheet(self.lineStyle_9)
        #  self.Button.setStyleSheet(self.buttonStyle_9)
        #  ------------------[ Coloring Kit ]---------------------



        self.editmode = 0 # Edit mode is off, 1 = on
        self.font_size_text_edit = 14
        self.font_family_text_edit = "Monospace"

        # Create the main vertical layout
        main_vertical_layout = QVBoxLayout(self.dialog)

        # Create the first horizontal layout with QTreeView and QTextEdit
        top_horizontal_control_layout = QHBoxLayout()
        top_horizontal_layout = QHBoxLayout()

        #style = QApplication.style()
        #self.bucket_icon = style.standardIcon(QStyle.SP_DriveNetIcon)
        #self.folder_icon = style.standardIcon(QStyle.SP_DirIcon)
        #self.file_icon = style.standardIcon(QStyle.SP_FileIcon)

        # bucket view comes first in the UI and disappears after a bucket is chosen.
        self.s3_bucket_view = QListWidget()
        self.s3_bucket_view.setFixedWidth(250) 
        self.s3_bucket_view.setStyleSheet("QListWidget { background-color: #f0f0f0; border: 1px solid #C06000; font-size: 16px; }")
        self.s3_bucket_view.setVisible(True)
        self.s3_bucket_view.clicked.connect(self.on_bucket_clicked)
        top_horizontal_layout.addWidget(self.s3_bucket_view)

        # Protection against viewing protected buckets
        self.s3_bucket_safety = ['isengard','do-not','delete','cloudtrail','aws','log','logs','audit','secure','pii','hipaa','gpdr','awsu','aws-tc']
        # Protection against deleting protected filenames
        self.s3_file_safety = ['isengard','do-not','delete','cloudtrail','aws','logs','log','awsu','aws-tc','config','credentials','audit']
        # Bypass protections
        # self.s3_bucket_safety = ['']
        # self.s3_file_safety = ['']

        # self.model = QStandardItemModel()
        self.model = CustomStandardItemModel() # Customized to prevent editing of names in the viewer
        self.s3_tree_view = QTreeView()
        self.s3_tree_view.setFixedWidth(300) 
        self.s3_tree_view.setStyleSheet("QTreeView { background-color: #f0f0f0; border: 1px solid #C06000; }")

        self.s3_tree_view.clicked.connect(self.on_item_clicked)
        self.s3_tree_view.setModel(self.model)
        self.model.setHorizontalHeaderLabels(['S3 Objects'])

        self.s3_text_edit = QTextEdit()
        # self.s3_text_edit.setStyleSheet("QTextEdit { background-color: #f0f0f0; border: 1px solid #C06000; font-size: 14px; }")
        self.text_edit_gray()
        self.s3_text_edit.setReadOnly(True) # Read only
        #self.s3_text_edit.setReadOnly(False) # Editable
        top_horizontal_layout.addWidget(self.s3_tree_view)
        top_horizontal_layout.addWidget(self.s3_text_edit)
        main_vertical_layout.addLayout(top_horizontal_control_layout)
        main_vertical_layout.addLayout(top_horizontal_layout)

        self.s3_list_view = QListWidget()
        self.s3_list_view.setStyleSheet("QListWidget { background-color: #f0f0f0; border: 1px solid #C06000; }")
        self.s3_list_view.setVisible(False)
        top_horizontal_layout.addWidget(self.s3_list_view)
        
        self.s3_metadata_view = QListWidget()
        self.s3_metadata_view.setStyleSheet("QListWidget { background-color: #f0f0f0; border: 1px solid #C06000; }")
        self.s3_metadata_view.setVisible(False)
        top_horizontal_layout.addWidget(self.s3_metadata_view)

        # Create the second horizontal layout for the buttons
        bottom_horizontal_layout = QHBoxLayout()
        left_vertical_layout = QVBoxLayout()
        right_vertical_layout = QVBoxLayout()
        bottom_area_layout = QVBoxLayout()
        current_area_layout = QVBoxLayout()
        current_area_layout.setSpacing(5)


        # Displays the current bucket, folder path, and file
        clabel_1 = QLabel("Bucket:")
        clabel_1.setFixedWidth(80)
        clabel_2 = QLabel("Folder-path:")
        clabel_2.setFixedWidth(80)
        clabel_3 = QLabel("File:")
        clabel_3.setFixedWidth(80)
        self.current_bucket = QLineEdit()
        self.current_bucket.setReadOnly(True)
        self.current_bucket.setStyleSheet(self.lineStyle_9)
        self.current_folder_path = QLineEdit()
        self.current_folder_path.setReadOnly(True)
        self.current_folder_path.setStyleSheet(self.lineStyle_9)
        self.current_file = QLineEdit()
        self.current_file.setReadOnly(True)
        self.current_file.setStyleSheet(self.lineStyle_9)

        clabel_4a = QLabel("Modified:")
        clabel_4a.setFixedWidth(80)
        clabel_4b = QLabel("   Length:")
        clabel_4b.setFixedWidth(60)
        clabel_4c = QLabel("   Type:")
        clabel_4c.setFixedWidth(60)

        clabel_5a = QLabel("Version:")
        clabel_5a.setFixedWidth(80)
        clabel_5b = QLabel("Metadata:")
        clabel_5b.setFixedWidth(80)
        self.current_last_modified = QLineEdit()
        self.current_last_modified.setFixedWidth(140)
        self.current_last_modified.setReadOnly(True)
        self.current_last_modified.setStyleSheet(self.lineStyle_9)
        self.current_content_length = QLineEdit()
        self.current_content_length.setFixedWidth(100)
        self.current_content_length.setReadOnly(True)
        self.current_content_length.setStyleSheet(self.lineStyle_9)
        self.current_version_id = QLineEdit()
        self.current_version_id.setReadOnly(True)
        self.current_version_id.setFixedWidth(260)
        self.current_version_id.setStyleSheet(self.lineStyle_9)
        self.current_content_type = QLineEdit()
        self.current_content_type.setReadOnly(True)
        self.current_content_type.setStyleSheet(self.lineStyle_9)
        self.current_metadata_str = QLineEdit()
        self.current_metadata_str.setReadOnly(True)
        self.current_metadata_str.setStyleSheet(self.lineStyle_9)

        current_row_1 = QHBoxLayout()
        current_row_2 = QHBoxLayout()
        current_row_3 = QHBoxLayout()
        current_row_4 = QHBoxLayout()
        current_row_5 = QHBoxLayout() 

        current_area_layout.addLayout(current_row_1)
        current_area_layout.addLayout(current_row_2)
        current_area_layout.addLayout(current_row_3)
        current_area_layout.addLayout(current_row_4)
        current_area_layout.addLayout(current_row_5)
        main_vertical_layout.addLayout(current_area_layout)
        current_row_1.addWidget(clabel_1)
        current_row_1.addWidget(self.current_bucket)
        current_row_2.addWidget(clabel_2)
        current_row_2.addWidget(self.current_folder_path)
        current_row_3.addWidget(clabel_3)
        current_row_3.addWidget(self.current_file)

        current_row_4.addWidget(clabel_4a)
        current_row_4.addWidget(self.current_last_modified)
        current_row_4.addWidget(clabel_4b)
        current_row_4.addWidget(self.current_content_length)
        current_row_4.addWidget(clabel_4c)
        current_row_4.addWidget(self.current_content_type)

        current_row_5.addWidget(clabel_5a)
        current_row_5.addWidget(self.current_version_id)
        current_row_5.addWidget(clabel_5b)
        current_row_5.addWidget(self.current_metadata_str)

        bSpace = 3
        button_row_1 = QHBoxLayout()
        button_row_1.setSpacing(bSpace)
        button_r1_c1 = QVBoxLayout()
        button_r1_c1.addStretch
        button_r1_c2 = QVBoxLayout()
        button_r1_c2.addStretch
        button_r1_c3 = QVBoxLayout()
        button_r1_c3.addStretch
        button_r1_c4 = QVBoxLayout()
        button_r1_c5 = QVBoxLayout()
        button_r1_c6 = QVBoxLayout()
        button_r1_c7 = QVBoxLayout()

        button_r1_c8 = QVBoxLayout()
        button_r1_c8.setSpacing(5)
        button_row_1.addLayout(button_r1_c1)
        button_row_1.addLayout(button_r1_c2)
        button_row_1.addLayout(button_r1_c3)
        button_row_1.addStretch()
        button_row_1.addLayout(button_r1_c4)
        button_row_1.addLayout(button_r1_c5)
        button_row_1.addLayout(button_r1_c6)
        button_row_1.addLayout(button_r1_c7)
        button_row_1.addStretch()
        #button_row_1.addLayout(button_r1_c8)
        #button_row_1.addStretch()

        button_row_2 = QHBoxLayout()
        bottom_area_layout.addLayout(button_row_1)
        bottom_area_layout.addLayout(button_row_2)
        main_vertical_layout.addLayout(bottom_area_layout)

        self.buttonStyle_2 = """
        QPushButton {
            background-color: #FFFFCC;
            color: #000000;
            font-family: Arial; 
            font-size: 12px;    
            font-weight: bold;  
            font-style: normal;  
            border: 2px solid #222222;
            border-radius: 0px;
            }
            QPushButton:hover { background-color: #FFC200; }
            QPushButton:pressed { background-color: #000000; color: #FFFFFF; }
        """

        self.buttonStyle_2g = """
        QPushButton {
            background-color: #E6E6E6;
            color: #000000;
            font-family: Arial; 
            font-size: 12px;    
            font-weight: bold;  
            font-style: normal;  
            border: 2px solid #C2C2C2;
            border-radius: 0px;
            }
            QPushButton:hover { background-color: #3a3a3a; color: #FFFFFF;}
            QPushButton:pressed { background-color: #FF0000; color: #000000; }
        """


        self.buttonStyle_5 = """
        QPushButton {
            background-color: #E6E6E6;
            color: #000000;
            font-family: Arial; 
            font-size: 12px;    
            font-weight: normal;  
            font-style: normal; 
            border: 2px solid #C2C2C2;
            border-radius: 7px;
            }
            QPushButton:hover { background-color: #3a3a3a; color: #FFFFFF;}
            QPushButton:pressed { background-color: #FF0000; color: #000000; }
        """

        self.buttonStyle_6 = """
        QPushButton {
            background-color: #C4E0EF;
            color: #000000;
            font-family: Arial; 
            font-size: 12px;    
            font-weight: bold;  
            font-style: normal; 
            border: 2px solid #000000;
            border-radius: 7px;
            }
            QPushButton:hover { background-color: #3a3a3a; color: #FFFFFF;}
            QPushButton:pressed { background-color: #FF0000; color: #000000; }
        """

        # -- Edit radio buttons
        # Create the radio buttons
        self.radioButtonEditOn = QRadioButton("Edit On")
        self.radioButtonEditOff = QRadioButton("Edit Off")
        # Set default checked state
        self.radioButtonEditOff.setChecked(True)
        # Connect the signal to the slot
        self.radioButtonEditOn.toggled.connect(self.editModeChanged)
        self.radioButtonEditOff.toggled.connect(self.editModeChanged)
        top_horizontal_control_layout.addStretch()
        top_horizontal_control_layout.addWidget(self.radioButtonEditOn)
        top_horizontal_control_layout.addWidget(self.radioButtonEditOff)

        # -- Copy to Clipboard button for S3 Editor
        self.copyS3 = QPushButton("📋") 
        self.copyS3.setToolTip("Copy S3 editor to clipboard.")
        self.copyS3.setFixedSize(20,20)
        self.copyS3.setStyleSheet(self.buttonStyle_6)
        self.copyS3.clicked.connect(self.clipboard_text_edit)  ## copy to clipboard
        top_horizontal_control_layout.addWidget(self.copyS3)

        font = QFont()
        font.setItalic(True)

        bW = 120
        bH = 20
        ## -- buttons ---
        column_s3_view = QLabel("View Controls")
        column_s3_view.setFont(font)
        column_s3_view.setStyleSheet("color: #C06000;")
        column_s3_view_spacer = QLabel(" ")
        button_r1_c1.addWidget(column_s3_view)
        button_r1_c1.addWidget(column_s3_view_spacer)

        self.b_011 = QPushButton('Hide Buckets')
        self.b_011.clicked.connect(self.show_hide_bucket_panel)
        self.b_011.setToolTip("View all buckets in S3.")
        self.b_011.setFixedSize(bW,bH)
        self.b_011.setStyleSheet(self.buttonStyle_6)
        button_r1_c1.addWidget(self.b_011)

        self.b_015 = QPushButton('Hide Folders')
        self.b_015.clicked.connect(self.show_hide_folder_panel)
        self.b_015.setToolTip("")
        self.b_015.setFixedSize(bW,bH)
        self.b_015.setStyleSheet(self.buttonStyle_6)
        button_r1_c1.addWidget(self.b_015)

       # self.b_033 = QPushButton('Show Versions')
        self.b_033 = QPushButton(' ')
       # self.b_033.clicked.connect(self.show_hide_version_panel)
        self.b_033.setFixedSize(bW,bH)
       # self.b_033.setStyleSheet(self.buttonStyle_6)
        self.b_033.setStyleSheet(self.buttonStyle_5)
        button_r1_c1.addWidget(self.b_033)

     #   self.b_016 = QPushButton('Show Metadata')
        self.b_016 = QPushButton(' ')
       # self.b_016.clicked.connect(self.show_hide_metadata_panel)
        self.b_016.setFixedSize(bW,bH)
       # self.b_016.setStyleSheet(self.buttonStyle_6)
        self.b_016.setStyleSheet(self.buttonStyle_5)
        button_r1_c1.addWidget(self.b_016)

        ##-- Folders
        column_Folders = QLabel("S3 Folders")
        column_Folders.setFont(font)
        column_Folders.setStyleSheet("color: #C06000;")
        column_Folders_spacer = QLabel(" ")
        button_r1_c2.addWidget(column_Folders)
        button_r1_c2.addWidget(column_Folders_spacer)

        b_023 = QPushButton('New S3 Folder')
        b_023.clicked.connect(self.create_folder)
        b_023.setToolTip("")
        b_023.setFixedSize(bW,bH)
        b_023.setStyleSheet(self.buttonStyle_6)
        button_r1_c2.addWidget(b_023)

        b_012 = QPushButton('Expand All')
        b_012.clicked.connect(self.expand)
        b_012.setToolTip("View all items in bucket.")
        b_012.setFixedSize(bW,bH)
        b_012.setStyleSheet(self.buttonStyle_6)
        button_r1_c2.addWidget(b_012)

        b_013 = QPushButton('Contract All')
        b_013.clicked.connect(self.contract)
        b_013.setToolTip("")
        b_013.setFixedSize(bW,bH)
        b_013.setStyleSheet(self.buttonStyle_6)
        button_r1_c2.addWidget(b_013)

        b_014 = QPushButton('Refresh')
        b_014.clicked.connect(self.refresh)
        b_014.setToolTip("")
        b_014.setFixedSize(bW,bH)
        b_014.setStyleSheet(self.buttonStyle_6)
        button_r1_c2.addWidget(b_014)


        ## --- Files
        column_s3_file = QLabel("S3 Files")
        column_s3_file.setFont(font)
        column_s3_file.setStyleSheet("color: #C06000;")
        column_s3_file_spacer = QLabel("")
        button_r1_c3.addWidget(column_s3_file)
        button_r1_c3.addWidget(column_s3_file_spacer)

        b_032 = QPushButton('New S3 File')
        b_032.clicked.connect(self.new_file)
        b_032.setToolTip("")
        b_032.setFixedSize(bW,bH)
        b_032.setStyleSheet(self.buttonStyle_6)
        button_r1_c3.addWidget(b_032)

        b_021 = QPushButton('Upload Files to S3')
        b_021.clicked.connect(self.upload_files)
        b_021.setToolTip("")
        b_021.setFixedSize(bW,bH)
        b_021.setStyleSheet(self.buttonStyle_6)
        button_r1_c3.addWidget(b_021)

        b_022 = QPushButton('Download Files')
        b_022.clicked.connect(self.download_files)
        b_022.setToolTip("")
        b_022.setFixedSize(bW,bH)
        b_022.setStyleSheet(self.buttonStyle_6)
        button_r1_c3.addWidget(b_022)

        b_024 = QPushButton('Delete S3 Files')
        b_024.clicked.connect(self.delete)
        b_024.setToolTip("")
        b_024.setFixedSize(bW,bH)
        b_024.setStyleSheet(self.buttonStyle_6)
        button_r1_c3.addWidget(b_024)

        #-- Editor
        column_Command = QLabel("S3 File Editor")
        column_Command.setFont(font)
        column_Command.setStyleSheet("color: #C06000;")
        column_Command_spacer = QLabel("")
        button_r1_c4.addWidget(column_Command)
        button_r1_c4.addWidget(column_Command_spacer)
    
        # Save changes changes color when in editmode
        self.b_031 = QPushButton('Save Changes')
        self.b_031.clicked.connect(self.store_file)
        self.b_031.setToolTip("")
        self.b_031.setFixedSize(bW,bH)
        self.b_031.setStyleSheet(self.buttonStyle_2g)
        button_r1_c4.addWidget(self.b_031)

        self.b_017 = QPushButton('Interpret')
        # self.b_017 = QPushButton(' ')
        self.b_017.clicked.connect(self.interpret_file)
        self.b_017.setFixedSize(bW,bH)
        self.b_017.setStyleSheet(self.buttonStyle_2)
        button_r1_c4.addWidget(self.b_017)

        b_034 = QPushButton(' ')
       # b_034 = QPushButton('Update Metadata')
        b_034.clicked.connect(self.expand)
        b_034.setToolTip("")
        b_034.setFixedSize(bW,bH)
        b_034.setStyleSheet(self.buttonStyle_2g)
        button_r1_c4.addWidget(b_034)

        b_035 = QPushButton(' ')
        # b_035 = QPushButton('Restore Version')
        b_035.clicked.connect(self.expand)
        b_035.setToolTip("")
        b_035.setFixedSize(bW,bH)
        b_035.setStyleSheet(self.buttonStyle_2g)
        #button_r1_c4.addWidget(b_035)


# Load Editor
        column_Command = QLabel("LOAD Editor")
        column_Command.setFont(font)
        column_Command.setStyleSheet("color: #C06000;")
        column_Command_spacer = QLabel("")
        button_r1_c5.addWidget(column_Command)
        button_r1_c5.addWidget(column_Command_spacer)

        b_041 = QPushButton('COMMAND')
        b_041.clicked.connect(self.load_command)
        b_041.setToolTip("")
        b_041.setFixedSize(bW,bH)
        b_041.setStyleSheet(self.buttonStyle_2)
        button_r1_c5.addWidget(b_041)

        b_042 = QPushButton('INPUT')
        b_042.clicked.connect(self.load_input)
        b_042.setToolTip("")
        b_042.setFixedSize(bW,bH)
        b_042.setStyleSheet(self.buttonStyle_2)
        button_r1_c5.addWidget(b_042)

        b_043 = QPushButton('RESPONSE')
        b_043.clicked.connect(self.load_response)
        b_043.setToolTip("")
        b_043.setFixedSize(bW,bH)
        b_043.setStyleSheet(self.buttonStyle_2)
        button_r1_c5.addWidget(b_043)

        b_044 = QPushButton('')
        b_044.clicked.connect(self.expand)
        b_044.setToolTip("")
        b_044.setFixedSize(bW,bH)
        b_044.setStyleSheet(self.buttonStyle_2g)
        # button_r1_c5.addWidget(b_044)

# Store Editor
        column_Input = QLabel("STORE Editor")
        column_Input.setFont(font)
        column_Input.setStyleSheet("color: #C06000;")
        column_Input_spacer = QLabel("")
        button_r1_c6.addWidget(column_Input)
        button_r1_c6.addWidget(column_Input_spacer)

        b_051 = QPushButton('COMMAND')
        b_051.clicked.connect(self.store_command)
        b_051.setToolTip("")
        b_051.setFixedSize(bW,bH)
        b_051.setStyleSheet(self.buttonStyle_2)
        button_r1_c6.addWidget(b_051)

        b_052 = QPushButton('INPUT')
        b_052.clicked.connect(self.store_input)
        b_052.setToolTip("")
        b_052.setFixedSize(bW,bH)
        b_052.setStyleSheet(self.buttonStyle_2)
        button_r1_c6.addWidget(b_052)

        b_053 = QPushButton('RESPONSE')
        b_053.clicked.connect(self.store_response)
        b_053.setToolTip("")
        b_053.setFixedSize(bW,bH)
        b_053.setStyleSheet(self.buttonStyle_2)
        button_r1_c6.addWidget(b_053)

        b_054 = QPushButton(' ')
        b_054.clicked.connect(self.expand)
        b_054.setToolTip("")
        b_054.setFixedSize(bW,bH)
        b_054.setStyleSheet(self.buttonStyle_5)
        # button_r1_c6.addWidget(b_054)

# - Editor Controls 
        column_Response = QLabel("Editor Controls")
        column_Response.setFont(font)
        column_Response.setStyleSheet("color: #C06000;")
        column_Response_spacer = QLabel("")
        button_r1_c7.addWidget(column_Response)
        button_r1_c7.addWidget(column_Response_spacer)

        b_061 = QPushButton('Clear Editor')
        b_061.clicked.connect(self.clear_text_editor)
        b_061.setToolTip("")
        b_061.setFixedSize(bW,bH)
        b_061.setStyleSheet(self.buttonStyle_2)
        button_r1_c7.addWidget(b_061)

        b_062 = QPushButton('Increase Font')
        b_062.clicked.connect(self.increase_font_size)
        b_062.setToolTip("")
        b_062.setFixedSize(bW,bH)
        b_062.setStyleSheet(self.buttonStyle_2)
        button_r1_c7.addWidget(b_062)

        b_063 = QPushButton('Decrease Font')
        b_063.clicked.connect(self.decrease_font_size)
        b_063.setToolTip("")
        b_063.setFixedSize(bW,bH)
        b_063.setStyleSheet(self.buttonStyle_2)
        button_r1_c7.addWidget(b_063)

        b_064 = QPushButton('R1C6-4')
        b_064.clicked.connect(self.expand)
        b_064.setToolTip("")
        b_064.setFixedSize(bW,bH)
        b_064.setStyleSheet(self.buttonStyle_5)
        # button_r1_c7.addWidget(b_064)

        # Adding the OK button to the right vertical layout
        ok_button = QPushButton('OK')
        b_044.setToolTip("")
        b_044.setFixedSize(bW,bH)
        ok_button.setStyleSheet(self.buttonStyle_9)
        ok_button.clicked.connect(self.well_ok_then)
        button_row_2.addStretch()
        button_row_2.addWidget(ok_button)

        # Add vertical layouts to the bottom horizontal layout
        bottom_horizontal_layout.addLayout(left_vertical_layout)
        bottom_horizontal_layout.addLayout(right_vertical_layout)
        main_vertical_layout.addLayout(bottom_horizontal_layout)

        # Load from AWS 
        self.show_buckets()
        self.dialog.exec_()

# ------------- S3 Methods -------------------------------------

    def well_ok_then(self):
        self.dialog.close()

## ------------  Show/Hide panels ---------------------------
    def show_hide_bucket_panel(self):
        if self.s3_bucket_view.isVisible():
            self.s3_bucket_view.setVisible(False) # Hide it
            self.b_011.setText("Show Buckets")
        else:
            self.s3_bucket_view.setVisible(True) # Show it
            self.b_011.setText("Hide Buckets")

    def show_hide_folder_panel(self):
        if self.s3_tree_view.isVisible():
            self.s3_tree_view.setVisible(False) # Hide it
            self.b_015.setText("Show Folders")
        else:
            self.s3_tree_view.setVisible(True) # Show it
            self.b_015.setText("Hide Folders")

    def show_hide_version_panel(self):
        if self.s3_list_view.isVisible():
            self.s3_list_view.setVisible(False) # Hide it
            self.b_033.setText("Show Versions")
        else:
            self.s3_list_view.setVisible(True) # Show it
            self.b_033.setText("Hide Versions")

    def show_hide_metadata_panel(self):
        if self.s3_metadata_view.isVisible():
            self.s3_metadata_view.setVisible(False) # Hide it
            self.b_016.setText("Show Metadata")
        else:
            self.s3_metadata_view.setVisible(True) # Show it
            self.b_016.setText("Hide Metadata")

## ------------  Ascend to s3 I/O ---------------------------

    def store_command(self):
        text_to_store = self.s3_text_edit.toPlainText()
        self.edit_1.setPlainText(text_to_store)

    def store_input(self):
        text_to_store = self.s3_text_edit.toPlainText()
        self.edit_2.setPlainText(text_to_store)

    def store_response(self):
        text_to_store = self.s3_text_edit.toPlainText()
        self.edit_3.setPlainText(text_to_store)

    def load_command(self):
        text_to_load = self.edit_1.toPlainText()
        self.s3_text_edit.setPlainText(text_to_load)
        self.b_031.setStyleSheet(self.buttonStyle_2)
        self.editmodemode = 1
        # self.s3_text_edit.setStyleSheet("QTextEdit { background-color: #FFFFFF; border: 1px solid #C06000; }" )
        self.text_edit_white()
        self.s3_text_edit.setReadOnly(False) # Editable
        self.radioButtonEditOn.setChecked(True)

    def load_input(self):
        text_to_load = self.edit_2.toPlainText()
        self.s3_text_edit.setPlainText(text_to_load)
        self.b_031.setStyleSheet(self.buttonStyle_2)
        self.editmodemode = 1
        # self.s3_text_edit.setStyleSheet("QTextEdit { background-color: #FFFFFF; border: 1px solid #C06000; }" )
        self.text_edit_white()
        self.s3_text_edit.setReadOnly(False) # Editable
        self.radioButtonEditOn.setChecked(True)

    def load_response(self):
        text_to_load = self.edit_3.toPlainText()
        self.s3_text_edit.setPlainText(text_to_load)
        self.b_031.setStyleSheet(self.buttonStyle_2)
        self.editmodemode = 1
        # self.s3_text_edit.setStyleSheet("QTextEdit { background-color: #FFFFFF; border: 1px solid #C06000; }" )
        self.text_edit_white()
        self.s3_text_edit.setReadOnly(False) # Editable
        self.radioButtonEditOn.setChecked(True)

    def clear_text_editor(self):
        self.s3_text_edit.clear()
        default_format = QTextCharFormat()
        self.s3_text_edit.setCurrentCharFormat(default_format)

    def increase_font_size(self):
        stylesheet = self.s3_text_edit.styleSheet()
        lines = stylesheet.split(';')
        new_lines = []
        font_size_set = False
        for line in lines:
            if 'font-size' in line:
                # Extract the current font size and increase it
                parts = line.split(':')
                if len(parts) > 1:
                    size_part = parts[1].strip().replace('px', '')
                    try:
                        current_size = int(size_part)
                        new_size = current_size + 2
                        new_line = f"font-size: {new_size}px"
                        font_size_set = True
                    except ValueError:
                        new_line = line.strip()
                else:
                    new_line = line.strip()
            else:
                new_line = line.strip()
            if new_line:
                new_lines.append(new_line)
        if not font_size_set:
            new_lines.append(f"font-size: {14 + 2}px")  # default to 14px + increment if not found
        new_stylesheet = '; '.join(new_lines)
        self.s3_text_edit.setStyleSheet(new_stylesheet)

    def decrease_font_size(self):
        stylesheet = self.s3_text_edit.styleSheet()
        lines = stylesheet.split(';')
        new_lines = []
        font_size_set = False
        for line in lines:
            if 'font-size' in line:
                # Extract the current font size and decrease it
                parts = line.split(':')
                if len(parts) > 1:
                    size_part = parts[1].strip().replace('px', '')
                    try:
                        current_size = int(size_part)
                        new_size = max(current_size - 2, 4)  # Decrease by 2 but not below 4px
                        new_line = f"font-size: {new_size}px"
                        font_size_set = True
                    except ValueError:
                        new_line = line.strip()
                else:
                    new_line = line.strip()
            else:
                new_line = line.strip()
            if new_line:
                new_lines.append(new_line)
        if not font_size_set:
            new_lines.append("font-size: 14px")  # Ensuring a default setting if font size was not previously set
        new_stylesheet = '; '.join(new_lines)
        self.s3_text_edit.setStyleSheet(new_stylesheet)

    def text_edit_white(self):
        # Retrieve the current stylesheet
        stylesheet = self.s3_text_edit.styleSheet()
        # Modify the background color to white
        new_stylesheet = self.update_background_color(stylesheet, '#FFFFFF')
        # Apply the updated stylesheet
        self.s3_text_edit.setStyleSheet(new_stylesheet)

    def text_edit_gray(self):
        # Retrieve the current stylesheet
        stylesheet = self.s3_text_edit.styleSheet()
        # Modify the background color to gray
        new_stylesheet = self.update_background_color(stylesheet, '#f0f0f0')
        # Apply the updated stylesheet
        self.s3_text_edit.setStyleSheet(new_stylesheet)

    def update_background_color(self, stylesheet, color):
        lines = stylesheet.split(';')
        new_lines = []
        color_set = False
        for line in lines:
            if 'background-color' in line:
                new_line = f"background-color: {color}"
                color_set = True
            else:
                new_line = line.strip()
            if new_line:
                new_lines.append(new_line)
        if not color_set:
            new_lines.append(f"background-color: {color}")  # Add background color if not present
        new_stylesheet = '; '.join(new_lines)
        return new_stylesheet


## --- active methods
    def add_objects(self, bucket_name, bucket_item):
        # List objects within the bucket
        paginator = self.clients['s3'].get_paginator('list_objects_v2')
        page_iterator = paginator.paginate(Bucket=bucket_name)

        ## all_objects is the hierarchical dictionary -- the virtual representation of the S3 hierarchy.
        ## These are added to the model when a click occurs to create the view that animates the tree display
        all_objects = {}
        for page in page_iterator:
            if 'Contents' in page:
                for obj in page['Contents']:
                    key = obj['Key']
                    parts = key.split('/') 
                    # The following code creates consistent folder handling in the node structure                  
                    if key.endswith('/'):
                        # Folder -- replace ending '/' except for empty strings in the list.
                        # I am going to later prune out the empty ones.
                        parts = [item + '/' if item else item for item in parts]
                    else:
                        # File -- replace ending '/' except for the final item which is the filename.
                        parts = [item + '/' if item and index != len(parts) - 1 else item for index, item in enumerate(parts)]
                    node = all_objects
                    for part in parts[:-1]:
                        if part not in node:
                            node[part] = {}
                        node = node[part]
                    if parts[-1] not in node:
                        node[parts[-1]] = None 
        

        def add_items(parent, nodes):
            for name, subnode in sorted(nodes.items()):
                if name == "" and (subnode is None or not subnode):  # Skip empty text nodes with no children
                    continue
                item = QStandardItem(name)
                parent.appendRow(item)
                if subnode is not None and subnode:  # Ensure there are actual subnodes to add
                    add_items(item, subnode)

        add_items(bucket_item, all_objects)

    def on_item_clicked(self, index):
        self.s3_text_edit.setText('')
        # Handle the case of an empty bucket.
        item = self.model.itemFromIndex(index)
         # Check if the item is a root item (no parent)
        if not item.parent():
            # This item is a root item, potentially a bucket.
            # Now check if this bucket is empty
            if self.model.rowCount(item.index()) == 0:
                print("The bucket is empty.")
                bucket_name = item.text()
                self.current_bucket.setText(bucket_name)
                self.current_folder_path.setText('')  # Clear out folder path
                self.current_file.setText('')         # Clear out file name
                return

        ####
        item = self.model.itemFromIndex(index)
        
        # User clicked on the bucket name
        if not item.parent(): #
            bucket_name = item.text()
            self.current_bucket.setText(bucket_name)
        self.current_folder_path.setText('')  # Clear out folder path
        self.current_file.setText('')         # Clear out file name


        #if item and self.s3_tree_view.model().rowCount(item.index()) == 0:  # Indicates it's a file
        if item:
            # Traverse up the tree to construct the full key path
            key_parts = []
            current_item = item

            # Traverse until you reach the root item
            while current_item.parent():
                key_parts.insert(0, current_item.text())
                current_item = current_item.parent()
            
            # Do not include the root item's text (bucket name) in the key parts
            bucket_name = current_item.text()

            self.current_bucket.setText(bucket_name)
            self.current_folder_path.setText('')
            self.current_file.setText('')
            ### print("DEBUG-01 bucket_name",bucket_name)

            #
            if key_parts: # the list is not empty
                if key_parts[-1]: # The last item is not empty
                    if key_parts[-1].endswith('/'):
                        # Folder path only
                        folder_path = ''.join(key_parts)
                        self.current_folder_path.setText(folder_path)
                        self.current_file.setText("")
                    else:
                        if len(key_parts) == 1: # No folder path, just file in root of bucket
                            folder_path=""
                            filename = key_parts[-1]
                            self.current_folder_path.setText(folder_path)
                            self.current_file.setText(filename)
                        else:
                            # Folder path and file
                            folder_path = ''.join(key_parts[:-1])
                            filename = key_parts[-1]
                            self.current_folder_path.setText(folder_path)
                            self.current_file.setText(filename)

            # key_parts.insert(0, current_item.text())  # Ensures the bucket name is part of the key_parts
            # I correctly formatted final slashes on ingress to avoid post slash processing
            key = ''.join(key_parts)  # Join parts to form the full key

            # Only try to fetch the object IF the user clicked on a file or folder.
            # If the user clicked on a bucket -- skip fetching from S3 because there is no object specified.
            ###  print("DEBUG-02",item," > ", self.s3_tree_view.model().rowCount(item.index()))
            if item and self.s3_tree_view.model().rowCount(item.index()) == 0:
                # Fetch the object from AWS S3 using the constructed key
                response = self.clients['s3'].get_object(Bucket=bucket_name, Key=key)
                content = response['Body'].read()
                # data bits
                last_modified = response['LastModified']
                last_modified_str = last_modified.strftime('%Y-%m-%d %H:%M:%S')
                content_type = response['ContentType']
                content_length = response['ContentLength']
                content_length_str = str(content_length)
                version_id = response.get('VersionId', None) 
                metadata = response.get('Metadata', {})
                metadata_str = ', '.join(f"{key}: {value}" for key, value in metadata.items())
                self.current_last_modified.setText(last_modified_str)
                self.current_content_type.setText(content_type)
                self.current_content_length.setText(content_length_str)
                self.current_version_id.setText(version_id)
                if version_id != None:
                    self.version() # populate version panel
                    #self.b_033.setStyleSheet(self.buttonStyle_5)  # Set Version control button to gray
                    #self.s3_list_view.setVisible(False)           # Close version control panel if it is open
                #else:
                    #self.b_033.setStyleSheet(self.buttonStyle_6)
                    #self.version() # populate version panel
                self.current_metadata_str.setText(metadata_str)
                #
                try:
                    # Attempt to decode the content as UTF-8 or show it as a raw representation on failure
                    self.s3_text_edit.setText(content.decode('utf-8'))
                except UnicodeDecodeError:
                    self.s3_text_edit.setText(repr(content))
            else:
                self.current_last_modified.setText("")
                self.current_content_type.setText("")
                self.current_content_length.setText("")
                self.current_version_id.setText("")
                self.b_033.setStyleSheet(self.buttonStyle_5)
                self.current_metadata_str.setText("")
            

    def get_current_bucket(self):
        # Retrieve the currently selected bucket name
        current_index = self.s3_tree_view.currentIndex()
        if not current_index.isValid():
            return None  # No selection

        current_item = self.model.itemFromIndex(current_index)
        while current_item.parent():
            current_item = current_item.parent()

        return current_item.text()

    # Clears the model, thereby clearing the tree view
    def clear_tree(self):
        self.model.clear()
        self.model.setHorizontalHeaderLabels(['S3 Buckets / Objects'])

    # -- clear empty lines from the model
    def prune_model(self):
        # Start pruning from the root item
        self._prune_recursive(self.model.invisibleRootItem())

    def _prune_recursive(self, item):
        print(self.model_to_json())
        # Iterate backwards through the children to safely remove items
        row_count = item.rowCount()
        for row in reversed(range(row_count)):
            child = item.child(row)
            if child.rowCount() == 0 and (child.text() == "" or child.text() is None):
                # If the item is empty (no children and no text), remove it
                item.removeRow(row)
            else:
                # Recursively prune child items
                self._prune_recursive(child)
    # -- clear empty lines from the model



    def expand(self):
        # Expands all nodes in the QTreeView
        self.s3_tree_view.expandAll()

    def contract(self):
        # Collapses all nodes in the QTreeView except the top level (buckets)
        self.s3_tree_view.collapseAll()
        # Iterate over each top-level item and ensure it is visible but not expanded
        root = self.model.invisibleRootItem()
        for i in range(root.rowCount()):
            index = self.model.index(i, 0)
            self.s3_tree_view.setExpanded(index, False)

    # Refreshes from S3 in case something has changed externally
    def refresh(self):
        self.on_bucket_clicked()


    def delete(self):
        # 
        # Prepare verify delete message
        msg_box = QMessageBox()
        msg_box.setIcon(QMessageBox.Warning)        # Set the icon of the message box
        msg_box.setWindowTitle("Confirm delete")    # Title of the message box
        msg_box.setText("Are you sure you want to delete this?")  # Main message text
        msg_box.setStandardButtons(QMessageBox.Ok | QMessageBox.Cancel)  # Include OK and Cancel buttons
        msg_box.setDefaultButton(QMessageBox.Cancel)  # Set the default focus on the Cancel button

        bucket=self.current_bucket.text()
        path=self.current_folder_path.text()
        filename=self.current_file.text()
        if filename:
            key=path+filename
            ## Verify it is not on the protected list 
            verify_key_safety = key.lower()
            for safety in self.s3_file_safety:
                    if safety in verify_key_safety:
                        QMessageBox.critical(None, "Error", f"Cannot delete. Name matches protected list: {str(key)}")
                        return
            verify_delete = msg_box.exec_()
            if verify_delete == QMessageBox.Ok:
                response = self.clients['s3'].delete_object(Bucket=bucket, Key=key)
                # print("Deleted message:",response)
                self.current_file.setText('')
                self.refresh()
        else:
            print("Directory delete")
        #response = self.clients['s3'].delete_object(Bucket=self.current_bucket, Key=key)
            objects_to_delete = self.clients['s3'].list_objects_v2(Bucket=bucket, Prefix=path)
            delete_keys = [{'Key': obj['Key']} for obj in objects_to_delete.get('Contents', [])]
            for safety_key in delete_keys:
                verify_key_safety = safety_key['Key'].lower()
                for safety in self.s3_file_safety:
                        if safety in verify_key_safety:
                            QMessageBox.critical(None, "Error", f"Cannot delete. Name matches protected list: {str(safety_key['Key'])}")
                            return
        
            ## Verify no files in the directory are on the protected list.
            # self.s3_file_safety
            ##
            if delete_keys:
                # print("DEBUG-Delete",delete_keys)
                verify_delete = msg_box.exec_()
                if verify_delete == QMessageBox.Ok:
                    # print("Delete confirmed.")
                    response = self.clients['s3'].delete_objects(
                        Bucket=bucket,
                        Delete={'Objects': delete_keys}
                        )
                    # print("Debug-After-Delete",response)
                    self.current_file.setText('')
                    self.current_folder_path.setText('')
                    self.refresh()
                else:
                    print("Delete cancelled.")
            else:
                print("DEBUG- No Objects to Delete")
                
            
    def upload_files(self):
        bucket=self.current_bucket.text()
        path=self.current_folder_path.text()
        localfiles, _ = QFileDialog.getOpenFileNames(None, "Select files to upload")
        for pathname in localfiles:
            filename = os.path.basename(pathname) 
            print(filename)
            key = path + filename
            try:
                response = self.clients['s3'].upload_file(pathname, bucket, key)
                print(f'Successfully uploaded {filename}')
                self.current_file.setText('')
                self.refresh()
            except Exception as e:
                print(f'Failed to upload {filename}: {str(e)}')


    def create_folder(self):
        # This function is called when the 'Create S3 Folder' button is pressed
        bucket_name = self.current_bucket.text()

        if not bucket_name.strip():  # .strip() removes any surrounding whitespace
            msg = QMessageBox()
            msg.setIcon(QMessageBox.Warning)
            msg.setText("No bucket selected")
            msg.setInformativeText("Please select a bucket before proceeding.")
            msg.setWindowTitle("Warning")
            msg.setStandardButtons(QMessageBox.Ok)
            msg.exec_()
            return

        path_name = self.current_folder_path.text()
        if bucket_name:  # Ensure there is a selected bucket
            folder_name, ok = QInputDialog.getText(self.dialog, 'Enter New Folder Name', 'Folder Name:')
            if ok and folder_name:  # Check if user pressed OK and input is not empty
                # folder_key = f'{folder_name.strip()}/'  # Ensure the key ends with a slash
                folder_key = f'{path_name.rstrip("/")}/{folder_name.strip()}/'
                if folder_key.startswith("/"):  # Remove leading slash if this is in the root
                    folder_key = folder_key.lstrip("/")
                ## Verify folder name is not protected
                verify_key_safety = folder_key.lower()
                for safety in self.s3_file_safety:
                    if safety in verify_key_safety:
                        QMessageBox.critical(None, "Error", f"Cannot create folder. Name matches protected list: {str(folder_key)}")
                        return
                self.clients['s3'].put_object(Bucket=bucket_name, Key=folder_key)
                self.refresh()  # Refresh the tree view
            else:
                print("Folder creation canceled or invalid folder name.")

    def print_json(self):
        print("\n>")
        print("MODEL=",self.model)
        print("JSON=",self.model_to_json())


    def model_to_json(self):
        # Convert the model to a JSON string
        root_item = self.model.invisibleRootItem()
        data_structure = self._extract_items(root_item)
        return json.dumps(data_structure, indent=4)  # Pretty print the JSON

    def _extract_items(self, item):
        # Recursively extract items from the model and build a nested list/dictionary
        result = []
        for row in range(item.rowCount()):
            child = item.child(row)
            node = {
                'text': child.text(),
                'children': self._extract_items(child)
            }
            result.append(node)
        return result

    def new_file(self):
        # This function is called when 'New File' button is pressed
        bucket_name = self.current_bucket.text()
        #try:
        #    verified = s3_client.get_bucket_versioning(Bucket=bucket_name)
        #    versioning_status = verified.get('Status', 'Not Enabled')
        #    QMessageBox.information(None, "Versioned Bucket.", "The versioning....")
        #    return versioning_status
        #except Exception as e:
        #    print(f"Error checking versioning status: {str(e)}")
        #    return None


        if not bucket_name.strip():  # .strip() removes any surrounding whitespace
            msg = QMessageBox()
            msg.setIcon(QMessageBox.Warning)
            msg.setText("No bucket selected")
            msg.setInformativeText("Please select a bucket before proceeding.")
            msg.setWindowTitle("Warning")
            msg.setStandardButtons(QMessageBox.Ok)
            msg.exec_()
            return

        path_name = self.current_folder_path.text()
        if bucket_name:  # Ensure there is a selected bucket
            file_name, ok = QInputDialog.getText(self.dialog, 'Enter New File Name', 'File Name:')
            if '/' in file_name:
                QMessageBox.information(None, "Invalid file name.", "The file name is not valid.")
                return
            if '.' not in file_name:
                file_name = file_name + '.txt'
            if ok and file_name:  # Check if user pressed OK and input is not empty
                file_key = f'{path_name.rstrip("/")}/{file_name.strip()}'
                if file_key.startswith("/"):  # Remove leading slash if this is in the root
                    file_key = file_key.lstrip("/")
                ## Verify protected
                verify_key_safety = file_key.lower()
                for safety in self.s3_file_safety:
                    if safety in verify_key_safety:
                        QMessageBox.critical(None, "Error", f"Cannot create file. Name matches protected list: {str(file_key)}")
                        return
                # print("DEBUG-FILE-KEY",file_key)
                self.b_031.setStyleSheet(self.buttonStyle_2)
                self.editmodemode = 1
                # self.s3_text_edit.setStyleSheet("QTextEdit { background-color: #FFFFFF; border: 1px solid #C06000; }" )
                self.text_edit_white()
                self.s3_text_edit.setReadOnly(False) # Editable
                self.radioButtonEditOn.setChecked(True)
                # show filename in line
                self.current_file.setText(file_name)
                ## --- You only check if the file already exists in a non-versioned bucket.
                ##     Because in a versioned bucket -- if it already exists -- the service is just going to create 
                ##     a new version. Which means that "head_object" method doesn't behave the asme.
                ##     so we need to skip the check if this is a versioned bucket.
                try:
                    response = self.clients['s3'].get_bucket_versioning(Bucket=bucket_name)
                    versioning_status = response.get('Status', 'Not Enabled')
                except Exception as e:
                    print(f"Error checking versioning status: {str(e)}")
                    versioning_status = None

                if versioning_status == 'Enabled':
                    # Versioned bucket: Skip checking if the object exists using head_object
                    # Skip straight to writing the file
                    pass
                else:
                    # Non-versioned bucket: Check if the object exists
                    try:
                        self.clients['s3'].head_object(Bucket=bucket_name, Key=file_key)
                        exists = True
                    except ClientError as e:
                        if e.response['Error']['Code'] == '404':
                            exists = False
                        else:
                            # Other errors might need to be handled differently
                            print(f"AWS Client Error: {e}")
                            # raise
                
                    if exists:
                        # Notify the user that the file exists
                        response = QMessageBox.question(None, "File Exists",
                                                    "The file already exists. Do you want to overwrite it?",
                                                    QMessageBox.Yes | QMessageBox.No, QMessageBox.No)
                        if response == QMessageBox.No:
                            return  # Stop the upload if the user does not want to overwrite

                # If the bucket is versioned or the file does not exist or the user wants to overwrite, proceed with upload
                self.clients['s3'].put_object(Bucket=bucket_name, Key=file_key)
                QMessageBox.information(None, "Upload Successful", "The file has been uploaded successfully.")
                #
                response = self.clients['s3'].get_object(Bucket=bucket_name, Key=file_key)
                content = response['Body'].read() 
                try:
                    # Attempt to decode the content as UTF-8 or show it as a raw representation on failure
                    self.s3_text_edit.setText(content.decode('utf-8'))
                except UnicodeDecodeError:
                    self.s3_text_edit.setText(repr(content))
                # data bits
                last_modified = response['LastModified']
                last_modified_str = last_modified.strftime('%Y-%m-%d %H:%M:%S')
                content_type = response['ContentType']
                content_length = response['ContentLength']
                content_length_str = str(content_length)
                version_id = response.get('VersionId', None) 
                metadata = response.get('Metadata', {})
                metadata_str = ', '.join(f"{key}: {value}" for key, value in metadata.items())
                self.current_last_modified.setText(last_modified_str)
                self.current_content_type.setText(content_type)
                self.current_content_length.setText(content_length_str)
                self.current_version_id.setText(version_id)
                self.current_metadata_str.setText(metadata_str)
                #
                self.refresh()  # Refresh the tree view
                self.find_and_select_item(bucket_name, path_name, file_name)
            else:
                print("File creation canceled or invalid file name.")


    def editModeChanged(self):
        if self.radioButtonEditOn.isChecked():
            # print("Edit mode is ON")
            # self.b_batch.setVisible(True)  # Show the button
            self.b_031.setStyleSheet(self.buttonStyle_2)
            self.editmodemode = 1
            self.text_edit_white()
            # self.s3_text_edit.setStyleSheet("QTextEdit { background-color: #FFFFFF; border: 1px solid #C06000; }" )
            self.s3_text_edit.setReadOnly(False) # Editable
        else:
            # print("Edit mode is OFF")
            # self.b_batch.setVisible(False)  # Hide the button
            self.b_031.setStyleSheet(self.buttonStyle_2g)
            self.editmode = 0
            # self.s3_text_edit.setStyleSheet("QTextEdit { background-color: #f0f0f0; border: 1px solid #C06000; }")
            self.text_edit_gray()
            self.s3_text_edit.setReadOnly(True) # Read only


    def find_and_select_item(self, bucket, path, filename):
        # Start by finding the bucket
        root = self.model.invisibleRootItem()
        bucket_item = None
        for i in range(root.rowCount()):
            item = root.child(i)
            if item.text() == bucket:
                bucket_item = item
                break

        if not bucket_item:
            print("Bucket not found")
            return

        # Now, navigate through the path to find the file
        path_parts = path.strip('/').split('/')  # Assume path is a string like "folder1/folder2/"
        current_item = bucket_item
        for part in path_parts:
            found = False
            for j in range(current_item.rowCount()):
                child = current_item.child(j)
                if child.text() == part + '/':  # Folders are marked with a trailing slash
                    current_item = child
                    found = True
                    break
            if not found:
                print("Path not found:", part)
                return

        # Finally, look for the file
        file_found = False
        for k in range(current_item.rowCount()):
            file_item = current_item.child(k)
            if file_item.text() == filename:  # Check for the file
                file_found = True
                # Select the item in the QTreeView
                index = file_item.index()
                self.s3_tree_view.selectionModel().select(index, QItemSelectionModel.ClearAndSelect)
                self.s3_tree_view.scrollTo(index)
                break

        if not file_found:
            print("File not found")

    def store_file(self):
        # Get the current text, bucket name, path, and filename
        if self.radioButtonEditOff.isChecked():
            QMessageBox.information(None, 'File not updated.', 'Edit mode is Off.')
            return

        text_to_store = self.s3_text_edit.toPlainText()
        bucket = self.current_bucket.text()
        path = self.current_folder_path.text().rstrip('/')  # Ensure no trailing slash
        filename = self.current_file.text()

        # Create the full S3 key for the object
        if path:
            s3_key = f"{path}/{filename}"
        else:
            s3_key = filename  # If no path is provided, use only filename

        verify_key_safety = s3_key.lower()
        for safety in self.s3_file_safety:
            if safety in verify_key_safety:
                QMessageBox.critical(None, "Error", f"Cannot overwrite file. Name matches protected list: {str(s3_key)}")
                return

        # Attempt to upload the file
        try:
            # Convert the string to bytes
            encoded_text = text_to_store.encode('utf-8')
            self.clients['s3'].put_object(Bucket=bucket, Key=s3_key, Body=encoded_text)
            QMessageBox.information(None, 'Success', 'File updated.')
        except Exception as e:
            QMessageBox.critical(None, 'Error', f'Failed to store file: {str(e)}')

## --- download files
    def choose_local_directory(self):
        return QFileDialog.getExistingDirectory(None, "Select Directory")

    def download_files(self):
        bucket = self.current_bucket.text()
        folder = self.current_folder_path.text()
        file = self.current_file.text()

        local_root_path = self.choose_local_directory()
        if not local_root_path:
            QMessageBox.warning(None, "No Directory Selected", "You must select a local directory to download files.")
            return
        
        if not bucket:
            QMessageBox.warning(None, "No Bucket Selected", "Please select a valid bucket.")
            return

        if file:
            # Downloading a single file
            print("DOWNLOAD ONE FILE")
            s3_object_key = os.path.join(folder, file)
            local_file_path = os.path.join(local_root_path, file)
            self.download_file(bucket, s3_object_key, local_file_path)
        else:
            # Downloading all files in the folder
            print("DOWNLOAD FOLDER")
            response = self.clients['s3'].list_objects_v2(Bucket=bucket, Prefix=folder)
            if 'Contents' not in response:
                QMessageBox.information(None, "No Files Found", "No files were found to download.")
                return
            
            for obj in response['Contents']:
                s3_object_key = obj['Key']
                if s3_object_key.endswith('/'):
                    continue  # Skip directories or empty keys

                local_file_path = os.path.join(local_root_path, s3_object_key.replace('/', os.sep))
                self.download_file(bucket, s3_object_key, local_file_path)

            QMessageBox.information(None, "Download Complete", "All files have been downloaded successfully.")

    def download_file(self, bucket, s3_object_key, local_file_path):
        local_file_directory = os.path.dirname(local_file_path)
        os.makedirs(local_file_directory, exist_ok=True)  # Ensure the directory exists

        # Download the file
        try:
            self.clients['s3'].download_file(bucket, s3_object_key, local_file_path)
            print(f'Downloaded {s3_object_key} to {local_file_path}')
        except Exception as e:
            QMessageBox.critical(None, "Download Failed", f"Failed to download the file: {str(e)}")

    def clipboard_text_edit(self):
        self.s3_text_edit.selectAll()
        self.s3_text_edit.copy()
        cursor = self.s3_text_edit.textCursor()
        cursor.clearSelection()
        self.s3_text_edit.setTextCursor(cursor)

    def version(self):
        bucket = self.current_bucket.text()
        folder = self.current_folder_path.text()
        file = self.current_file.text()

        if bucket == '': ## There is no bucket selected
            return 
        
        try:
            version_info = self.clients['s3'].get_bucket_versioning(Bucket=bucket)
            versioning_status = version_info.get('Status', 'Not Enabled')
        except Exception as e:
                print(f"Error checking versioning status: {str(e)}")
                versioning_status = None
        if versioning_status == 'Not Enabled':   ## The bucket does not have versioning enabled
            return
       
        if file == '':  ## Do you have a file selected? Files only.
            return
    
        self.s3_list_view.clear()

        paginator = self.clients['s3'].get_paginator('list_object_versions')
        try:
            for page in paginator.paginate(Bucket=bucket, Prefix=folder + file):
                for version in page.get('Versions', []):
                    if version['Key'] == folder + file:
                        version_id = version['VersionId']
                        last_modified = version['LastModified'].strftime('%Y-%m-%d %H:%M:%S')
                        size = version['Size']
                        self.s3_list_view.addItem(f"VERSION: {version_id}, MOD: {last_modified}, SIZE: {size} bytes")
        except Exception as e:
            QMessageBox.critical(None, "Error", f"An error occurred while fetching versions: {str(e)}")
            return



# ---
    def show_buckets(self):
        self.s3_bucket_view.clear()
        # Try to fetch the list of buckets and add them to the list widget
        try:
            response = self.clients['s3'].list_buckets()
            for bucket in response['Buckets']:
                self.s3_bucket_view.addItem(bucket['Name'])
        except Exception as e:
            print("Failed to list S3 buckets:", e)
    
    def on_bucket_clicked(self):
        # Set bucket
        bucket_name = self.s3_bucket_view.currentItem().text()
        verify_bucket_safety = bucket_name.lower()
        for safety in self.s3_bucket_safety:
            if safety in verify_bucket_safety:
                QMessageBox.critical(None, "Error", f"Name matches protected list: {str(bucket_name)}")
                return

        self.current_bucket.setText(bucket_name)
       # print("DEBUG-Bucket",self.current_bucket)
        self.model.clear()
        self.model.setHorizontalHeaderLabels(['S3 Buckets / Objects'])
        bucket_item = QStandardItem(bucket_name)
        self.model.appendRow(bucket_item)
        self.add_objects(bucket_name, bucket_item)
        self.expand()




#------------------------

## ------------  Interpret file ---------------------------

    def interpret_file(self):
        bucket = self.current_bucket.text()
        folder = self.current_folder_path.text()
        file = self.current_file.text()
        s3_object_key = os.path.join(folder, file)
        #local_file_path = os.path.join(local_root_path, file)

        _, file_extension = os.path.splitext(self.current_file.text())
        if file_extension not in ['.pptx', '.docx', '.xlsx', '.pdf','.md','.json','.csv','.html']:
            return
        # Create a temporary file and automatically handle its removal
        with tempfile.NamedTemporaryFile(suffix=file_extension,delete=False) as temp_file:
            temp_file_path = temp_file.name
            print("Temporary file path:", temp_file_path)
            # Download the file
            try:
                self.clients['s3'].download_file(bucket, s3_object_key, temp_file_path)
                print(f'Downloaded {s3_object_key} to {temp_file_path}')
            except Exception as e:
                QMessageBox.critical(None, "Download Failed", f"Failed to download the file: {str(e)}")
                return  # Exit the function if download fails
            # Ensure the file stream is flushed and file handle is closed
            temp_file.close()
            # Process the various file types
            # ------------------------------- Word Document
            if file_extension == '.docx':
                try:
                    doc = Document(temp_file_path)
                    fullText = []
                    for para in doc.paragraphs:
                        fullText.append(para.text)
                    text = '\n'.join(fullText)
                    self.s3_text_edit.setText(text)
                except Exception as e:
                    QMessageBox.critical(None, "DOCX Processing Failed", f"Failed to process the docx: {str(e)}")
            # ------------------------------- PDF Document
            elif file_extension == '.pdf':
                # Create a PDF document from the file
                try:
                    doc = fitz.open(temp_file_path)
                    text = ""
                    for page in doc:
                        text += page.get_text()
                    doc.close()
                    # Set the text in the QTextEdit
                    self.s3_text_edit.setText(text)
                except Exception as e:
                    QMessageBox.critical(None, "PDF Processing Failed", f"Failed to process the PDF: {str(e)}")
            # ------------------------------- Execel Spreadsheet (first worksheet only)
            elif file_extension == '.xlsx':
                try:
                    # Load the workbook and select the first sheet
                    wb = openpyxl.load_workbook(temp_file_path)
                    sheet = wb[wb.sheetnames[0]]
                    html = "<table border='1'>"
                    # Convert each row in the sheet to an HTML table row
                    for row in sheet.iter_rows(values_only=True):
                        html += "<tr>" + "".join(f"<td>{value}</td>" if value is not None else "<td></td>" for value in row) + "</tr>"
                    html += "</table>"
                    self.s3_text_edit.setHtml(html)
                except Exception as e:
                    QMessageBox.critical(None, "XLSX Processing Failed", f"Failed to process the xlsx: {str(e)}")
            elif file_extension == '.csv':
                csv_content = ""
                try:
                    with open(temp_file_path, mode='r', encoding='utf-8', errors='ignore') as file:
                        reader = csv.reader(file)
                        for row in reader:
                            # Convert each row to a string and append to csv_content
                            csv_content += ', '.join(row) + '\n'
                except FileNotFoundError:
                    print("The CSV file was not found.")
                except Exception as e:
                    print(f"An error occurred while reading the CSV file: {e}")
                self.s3_text_edit.setText(csv_content)
            elif file_extension == '.json':
                try:
                    with open(temp_file_path, 'r', encoding='utf-8', errors='ignore') as file:
                        content = json.load(file)
                        pretty_content = json.dumps(content, indent=4)
                        self.s3_text_edit.setText(pretty_content)
                except FileNotFoundError:
                    print("The JSON file was not found.")
                except Exception as e:
                    print(f"An error occurred while reading the JSON file: {e}")
            elif file_extension == '.html':
                html_content = ""
                try:
                    with open(temp_file_path, 'r', encoding='utf-8', errors='ignore') as file:
                        html_content = file.read()
                        self.s3_text_edit.setHtml(html_content)  # Use setHtml for HTML content
                except FileNotFoundError:
                    print("The HTML file was not found.")
                except Exception as e:
                    print(f"An error occurred while reading the HTML file: {e}")

            elif file_extension == '.md':
                try:
                    with open(temp_file_path, 'r', encoding='utf-8', errors='ignore') as file:
                        md_content = file.read()
                        html_content = markdown.markdown(md_content)
                except FileNotFoundError:
                    print("The Markdown file was not found.")
                except Exception as e:
                    print(f"An error occurred while reading the Markdown file: {e}")
                self.s3_text_edit.setHtml(html_content)
            elif file_extension == '.pptx':
                # self.process_pptx()
                ## ---
                try:
                    presentation = Presentation(temp_file_path)
                    all_text = []
                    for slide_number, slide in enumerate(presentation.slides, start=1):
                        slide_text = [f"Slide: {slide_number}"]
                        # Extract text from slide shapes
                        for shape in slide.shapes:  
                            if shape.has_text_frame:
                                for paragraph in shape.text_frame.paragraphs:
                                    for run in paragraph.runs:
                                        slide_text.append(run.text)
                        # Extract text from speaker notes
                        if slide.notes_slide:
                            notes = slide.notes_slide.notes_text_frame
                            if notes:
                                notes_text = ['Speaker Notes:']
                                for paragraph in notes.paragraphs:
                                    for run in paragraph.runs:
                                        notes_text.append(run.text)
                                slide_text.extend(notes_text)
                        all_text.extend(slide_text)
                    # Join all text and set it to edit_2
                    all_text_str = '\n'.join(all_text)
                    self.s3_text_edit.setPlainText(all_text_str)
                except Exception as e:
                    print(f"Failed to process PPTX: {str(e)}")  
    
                ## ---
                print("pptx")
            else:
                print("Unsupported file type.")
        
        os.remove(temp_file_path)  


    


