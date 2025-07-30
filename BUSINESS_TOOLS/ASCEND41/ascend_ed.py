# ascend_ed.py
# Forge Ascend v4.0
# Updated by Tom Stern, 05 NOV 2024
#
#   based on Ascend 1 -- first version January 22 2024 -- by Tom Stern

import sys
import os
from PyQt5.QtWidgets import (
    QApplication, QMainWindow, QWidget, QHBoxLayout, QVBoxLayout, QLabel, QPushButton, QFrame, QTextEdit, QLineEdit,
    QGridLayout, QGroupBox, QFormLayout, QFileDialog, QDialog, QMessageBox, QInputDialog, QCheckBox, QFileDialog
)
from PyQt5.QtNetwork import QNetworkAccessManager, QNetworkRequest, QNetworkReply
from PyQt5.QtGui import QTextCharFormat

from PyQt5.QtGui import QFont, QIcon
from PyQt5.QtCore import QUrl
from PyQt5.QtWidgets import QFileDialog
from pptx import Presentation
import shutil
import re
from docx import Document
import fitz  # Import PyMuPDF
import openpyxl
import csv
import json
import markdown
from bs4 import BeautifulSoup
from itertools import zip_longest

class FileMerge:
    def __init__(self, text_editor):

        print("begin")
        # self.text_editor = text_editor
        # self.filtered_files = {'1': [], '2': [], '3': [], '4': [],'5': [], '6': [],}

        # self.settings_list = []
        # default_values = {'check': False, 'directory': '', 'pattern': '', 'insert': ''}
        #for _ in range(11):
        #    self.settings_list.append(default_values.copy())   # Copy ensures each dictionary is separate
        # self.settings_list[0]['check'] = True

    def consolidate(self):
        self.dialog = QDialog()
        self.dialog.setWindowTitle("Consolidate Text Files")
        self.dialog.setMinimumWidth(1200)

        # ------------------[ Coloring Kit ]---------------------
        self.dialog.setStyleSheet("background-color: #E6E6E6; color: black;")
        self.buttonStyle_9 = """
        QPushButton {
            background-color: #C4E0EF;
            color: #000000;
            font-family: Arial; 
            font-size: 14px;    
            font-weight: normal;  
            font-style: normal; 
            border: 1px solid #000000;
            border-radius: 2px;
            }
            QPushButton:hover { background-color: #5b5b5b; color: #FFFFFF;}
            QPushButton:pressed { background-color: #FF0000; color: #000000; }
        """
        self.buttonStyle_9g = """
        QPushButton {
            background-color: #5F5F5F;
            color: #FFFFFF;
            font-family: Arial; 
            font-size: 14px;    
            font-weight: normal;  
            font-style: normal; 
            border: 1px solid #000000;
            border-radius: 2px;
            }
            QPushButton:hover { background-color: #5b5b5b; color: #FFFFFF;}
            QPushButton:pressed { background-color: #FF0000; color: #000000; }
        """ 
        self.textStyle_9 = """
        QTextEdit{
            background-color: #FFFFFF;
            color: #000000;
            font-family: Arial; 
            font-size: 16px;    
            font-weight: normal;  
            font-style: normal; 
            border: 1px solid #000000; }
        """
        self.lineStyle_9 = """
        QLineEdit{
            background-color: #FFFFFF;
            color: #000000;
            font-family: Arial; 
            font-size: 14px;    
            font-weight: normal;  
            font-style: normal; 
            border: 1px solid #000000; }
        """


        # ---- Layouts ----
        layout = QVBoxLayout()
        ed_layout = QHBoxLayout()
        ed_v1_layout = QVBoxLayout()
        ed_v2_layout = QVBoxLayout()
        ed_v3_layout = QVBoxLayout()
        ed_v4_layout = QVBoxLayout()
        ed_v5_layout = QVBoxLayout()
        ed_v6_layout = QVBoxLayout()
        ed_b1_layout = QVBoxLayout()
        ed_b2_layout = QVBoxLayout()
        ed_b3_layout = QVBoxLayout()
        ed_b4_layout = QVBoxLayout()
        ed_b5_layout = QVBoxLayout()
        ed_b6_layout = QVBoxLayout()
        allbutton_layout = QHBoxLayout()
        but_1_layout = QVBoxLayout()
        but_2_layout = QVBoxLayout()
        but_2_layout.addStretch()
        # ---- Assembly ---
        ed_layout.addLayout(ed_v1_layout)
        ed_layout.addLayout(ed_v2_layout)
        ed_layout.addLayout(ed_v3_layout)
        ed_layout.addLayout(ed_v4_layout)
        ed_layout.addLayout(ed_v5_layout)
        ed_layout.addLayout(ed_v6_layout)
        self.fl_1_label = QLabel("File List 1")
        self.fl_1 = QTextEdit()
        self.fl_1.setMinimumHeight(400)
        self.fl_1.setStyleSheet(self.textStyle_9)
        self.fl_2_label = QLabel("File List 2")
        self.fl_2 = QTextEdit()
        self.fl_2.setMinimumHeight(400)
        self.fl_2.setStyleSheet(self.textStyle_9)
        self.fl_3_label = QLabel("File List 3")
        self.fl_3 = QTextEdit()
        self.fl_3.setMinimumHeight(400)
        self.fl_3.setStyleSheet(self.textStyle_9)
        self.fl_4_label = QLabel("File List 4")
        self.fl_4 = QTextEdit()
        self.fl_4.setStyleSheet(self.textStyle_9)
        self.fl_4.setMinimumHeight(400)
        self.fl_5_label = QLabel("File List 5")
        self.fl_5 = QTextEdit()
        self.fl_5.setStyleSheet(self.textStyle_9)
        self.fl_5.setMinimumHeight(400)
        self.fl_6_label = QLabel("File List 6")
        self.fl_6 = QTextEdit()
        self.fl_6.setStyleSheet(self.textStyle_9)
        self.fl_6.setMinimumHeight(400)
        self.led_1 = QLineEdit()
        self.led_1.setStyleSheet(self.lineStyle_9)
        self.led_2 = QLineEdit()
        self.led_2.setStyleSheet(self.lineStyle_9)
        self.led_3 = QLineEdit()
        self.led_3.setStyleSheet(self.lineStyle_9)
        self.led_4 = QLineEdit()
        self.led_4.setStyleSheet(self.lineStyle_9)
        self.led_5 = QLineEdit()
        self.led_5.setStyleSheet(self.lineStyle_9)
        self.led_6 = QLineEdit()
        self.led_6.setStyleSheet(self.lineStyle_9)
        self.outdir_label = QLabel("Output Directory:")
        self.outdir = QLineEdit()
        self.outdir.setStyleSheet(self.lineStyle_9)
        self.suffix_label = QLabel("File Suffix:")
        self.suffix = QLineEdit()
        self.suffix.setStyleSheet(self.lineStyle_9)
        ed_v1_layout.addWidget(self.fl_1_label)
        ed_v1_layout.addWidget(self.fl_1)
        ed_v1_layout.addWidget(self.led_1)       
        ed_v1_layout.addLayout(ed_b1_layout)
        ed_v2_layout.addWidget(self.fl_2_label)
        ed_v2_layout.addWidget(self.fl_2)
        ed_v2_layout.addWidget(self.led_2)  
        ed_v2_layout.addLayout(ed_b2_layout)
        ed_v3_layout.addWidget(self.fl_3_label)
        ed_v3_layout.addWidget(self.fl_3)
        ed_v3_layout.addWidget(self.led_3)  
        ed_v3_layout.addLayout(ed_b3_layout)
        ed_v4_layout.addWidget(self.fl_4_label)
        ed_v4_layout.addWidget(self.fl_4)
        ed_v4_layout.addWidget(self.led_4)  
        ed_v4_layout.addLayout(ed_b4_layout)
        ed_v5_layout.addWidget(self.fl_5_label)
        ed_v5_layout.addWidget(self.fl_5)
        ed_v5_layout.addWidget(self.led_5)  
        ed_v5_layout.addLayout(ed_b5_layout)
        ed_v6_layout.addWidget(self.fl_6_label)
        ed_v6_layout.addWidget(self.fl_6)
        ed_v6_layout.addWidget(self.led_6)  
        ed_v6_layout.addLayout(ed_b6_layout)
        allbutton_layout.addLayout(but_1_layout)
        allbutton_layout.addStretch()
        allbutton_layout.addLayout(but_2_layout)
        layout.addLayout(ed_layout)
        layout.addLayout(allbutton_layout)

        self.vh_1 = QPushButton("Hide 1")
        self.vh_1.setFixedWidth(60)
        self.vh_1.setStyleSheet(self.buttonStyle_9)
        self.vh_1.clicked.connect(self.show_hide_1)
        self.vh_2 = QPushButton("Hide 2")
        self.vh_2.setFixedWidth(60)
        self.vh_2.setStyleSheet(self.buttonStyle_9)
        self.vh_2.clicked.connect(self.show_hide_2)
        self.vh_3 = QPushButton("Show 3")
        self.vh_3.setFixedWidth(60)
        self.vh_3.setStyleSheet(self.buttonStyle_9g)
        self.vh_3.clicked.connect(self.show_hide_3)
        self.vh_4 = QPushButton("Show 4")
        self.vh_4.setFixedWidth(60)
        self.vh_4.setStyleSheet(self.buttonStyle_9g)
        self.vh_4.clicked.connect(self.show_hide_4)
        self.vh_5 = QPushButton("Show 5")
        self.vh_5.setFixedWidth(60)
        self.vh_5.setStyleSheet(self.buttonStyle_9g)
        self.vh_5.clicked.connect(self.show_hide_5)
        self.vh_6 = QPushButton("Show 6")
        self.vh_6.setFixedWidth(60)
        self.vh_6.setStyleSheet(self.buttonStyle_9g)
        self.vh_6.clicked.connect(self.show_hide_6)
        but_1_layout.addSpacing(2)
        but_1_layout.addWidget(self.vh_1)
        but_1_layout.addWidget(self.vh_2)
        but_1_layout.addWidget(self.vh_3)
        but_1_layout.addWidget(self.vh_4)
        but_1_layout.addWidget(self.vh_5)
        but_1_layout.addWidget(self.vh_6)

        self.sel_1 = QPushButton("Select Files 1")
        self.sel_1.setFixedWidth(120)
        self.sel_1.setStyleSheet(self.buttonStyle_9)
        self.sel_1.clicked.connect(self.get_1)
        ed_b1_layout.addWidget(self.sel_1)

        self.sel_2 = QPushButton("Select Files 2")
        self.sel_2.setFixedWidth(120)
        self.sel_2.setStyleSheet(self.buttonStyle_9)
        self.sel_2.clicked.connect(self.get_2)
        ed_b2_layout.addWidget(self.sel_2)

        self.sel_3 = QPushButton("Select Files 3")
        self.sel_3.setFixedWidth(120)
        self.sel_3.setStyleSheet(self.buttonStyle_9)
        self.sel_3.clicked.connect(self.get_3)
        ed_b3_layout.addWidget(self.sel_3)

        self.sel_4 = QPushButton("Select Files 4")
        self.sel_4.setFixedWidth(120)
        self.sel_4.setStyleSheet(self.buttonStyle_9)
        self.sel_4.clicked.connect(self.get_4)
        ed_b4_layout.addWidget(self.sel_4)

        self.sel_5 = QPushButton("Select Files 5")
        self.sel_5.setFixedWidth(120)
        self.sel_5.setStyleSheet(self.buttonStyle_9)
        self.sel_5.clicked.connect(self.get_5)
        ed_b5_layout.addWidget(self.sel_5)

        self.sel_6 = QPushButton("Select Files 6")
        self.sel_6.setFixedWidth(120)
        self.sel_6.setStyleSheet(self.buttonStyle_9)
        self.sel_6.clicked.connect(self.get_6)
        ed_b6_layout.addWidget(self.sel_6)


        # Output and Merge control
        but_2_layout.addWidget(self.outdir_label)
        but_2_layout.addWidget(self.outdir)

        self.outdir_button = QPushButton("Select Output")
        self.outdir_button.setFixedWidth(120)
        self.outdir_button.setStyleSheet(self.buttonStyle_9)
        self.outdir_button.clicked.connect(self.get_outdir)
        but_2_layout.addWidget(self.outdir_button)

        but_2_layout.addWidget(self.suffix_label)
        but_2_layout.addWidget(self.suffix)

        but_2h_layout = QHBoxLayout()
        but_2h_layout.addStretch()
        but_2_layout.addLayout(but_2h_layout)
        self.c_button = QPushButton("Consolidate")
        self.c_button.setFixedWidth(120)
        self.c_button.setStyleSheet(self.buttonStyle_9)
        self.c_button.clicked.connect(self.merge_files)
        but_2h_layout.addWidget(self.c_button)
        self.ok_button = QPushButton("QUIT NO SAVE")
        self.ok_button.setFixedWidth(120)
        self.ok_button.setStyleSheet(self.buttonStyle_9)
        self.ok_button.clicked.connect(self.well_ok_then)
        but_2h_layout.addWidget(self.ok_button)

        # --- Set default invisibles ---
        self.fl_3_label.setVisible(False)
        self.sel_3.setVisible(False)
        self.fl_3.setVisible(False) # Hide it 
        self.led_3.setVisible(False) # Hide it      
        self.fl_4_label.setVisible(False)
        self.sel_4.setVisible(False)
        self.fl_4.setVisible(False) # Hide it
        self.led_4.setVisible(False) # Hide it 
        self.fl_5_label.setVisible(False)
        self.sel_5.setVisible(False)
        self.fl_5.setVisible(False) # Hide it 
        self.led_5.setVisible(False) # Hide it      
        self.fl_6_label.setVisible(False)
        self.sel_6.setVisible(False)
        self.fl_6.setVisible(False) # Hide it  
        self.led_6.setVisible(False) # Hide it  


        self.dialog.setLayout(layout)
        self.dialog.exec_()


    def show_hide_1(self):
        if self.fl_1.isVisible():
            self.fl_1_label.setVisible(False)
            self.sel_1.setVisible(False)
            self.fl_1.setVisible(False) # Hide it
            self.led_1.setVisible(False)  
            self.vh_1.setStyleSheet(self.buttonStyle_9g)
            self.vh_1.setText("Show 1")
        else:
            self.fl_1_label.setVisible(True)
            self.sel_1.setVisible(True)
            self.fl_1.setVisible(True) # Show it
            self.led_1.setVisible(True) 
            self.vh_1.setStyleSheet(self.buttonStyle_9)
            self.vh_1.setText("Hide 1")

    def show_hide_2(self):
        if self.fl_2.isVisible():
            self.fl_2_label.setVisible(False)
            self.sel_2.setVisible(False)
            self.fl_2.setVisible(False) # Hide it
            self.led_2.setVisible(False)  
            self.vh_2.setStyleSheet(self.buttonStyle_9g)
            self.vh_2.setText("Show 2")
        else:
            self.fl_2_label.setVisible(True)
            self.sel_2.setVisible(True)
            self.fl_2.setVisible(True) # Show it
            self.led_2.setVisible(True) 
            self.vh_2.setStyleSheet(self.buttonStyle_9)
            self.vh_2.setText("Hide 2")

    def show_hide_3(self):
        if self.fl_3.isVisible():
            self.fl_3_label.setVisible(False)
            self.sel_3.setVisible(False)
            self.fl_3.setVisible(False) # Hide it
            self.led_3.setVisible(False)  
            self.vh_3.setStyleSheet(self.buttonStyle_9g)
            self.vh_3.setText("Show 3")
        else:
            self.fl_3_label.setVisible(True)
            self.sel_3.setVisible(True)
            self.fl_3.setVisible(True) # Show it
            self.led_3.setVisible(True) 
            self.vh_3.setStyleSheet(self.buttonStyle_9)
            self.vh_3.setText("Hide 3")

    def show_hide_4(self):
        if self.fl_4.isVisible():
            self.fl_4_label.setVisible(False)
            self.sel_4.setVisible(False)
            self.fl_4.setVisible(False) # Hide it
            self.led_4.setVisible(False)  
            self.vh_4.setStyleSheet(self.buttonStyle_9g)
            self.vh_4.setText("Show 4")
        else:
            self.fl_4_label.setVisible(True)
            self.sel_4.setVisible(True)
            self.fl_4.setVisible(True) # Show it
            self.led_4.setVisible(True) 
            self.vh_4.setStyleSheet(self.buttonStyle_9)
            self.vh_4.setText("Hide 4")

    def show_hide_5(self):
        if self.fl_5.isVisible():
            self.fl_5_label.setVisible(False)
            self.sel_5.setVisible(False)
            self.fl_5.setVisible(False) # Hide it
            self.led_5.setVisible(False)  
            self.vh_5.setStyleSheet(self.buttonStyle_9g)
            self.vh_5.setText("Show 5")
        else:
            self.fl_5_label.setVisible(True)
            self.sel_5.setVisible(True)
            self.fl_5.setVisible(True) # Show it
            self.led_5.setVisible(True) 
            self.vh_5.setStyleSheet(self.buttonStyle_9)
            self.vh_5.setText("Hide 5")

    def show_hide_6(self):
        if self.fl_6.isVisible():
            self.fl_6_label.setVisible(False)
            self.sel_6.setVisible(False)
            self.fl_6.setVisible(False) # Hide it
            self.led_6.setVisible(False)  
            self.vh_6.setStyleSheet(self.buttonStyle_9g)
            self.vh_6.setText("Show 6")
        else:
            self.fl_6_label.setVisible(True)
            self.sel_6.setVisible(True)
            self.fl_6.setVisible(True) # Show it
            self.led_6.setVisible(True) 
            self.vh_6.setStyleSheet(self.buttonStyle_9)
            self.vh_6.setText("Hide 6")


    def get_1(self):
        options = QFileDialog.Options()
        file_paths, _ = QFileDialog.getOpenFileNames(None, "Select Files", "", "All Files (*)", options=options)
        if file_paths:
            self.fl_1.append("\n".join(file_paths))

    def get_2(self):
        options = QFileDialog.Options()
        file_paths, _ = QFileDialog.getOpenFileNames(None, "Select Files", "", "All Files (*)", options=options)
        if file_paths:
            self.fl_2.append("\n".join(file_paths))

    def get_3(self):
        options = QFileDialog.Options()
        file_paths, _ = QFileDialog.getOpenFileNames(None, "Select Files", "", "All Files (*)", options=options)
        if file_paths:
            self.fl_3.append("\n".join(file_paths))

    def get_4(self):
        options = QFileDialog.Options()
        file_paths, _ = QFileDialog.getOpenFileNames(None, "Select Files", "", "All Files (*)", options=options)
        if file_paths:
            self.fl_4.append("\n".join(file_paths))

    def get_5(self):
        options = QFileDialog.Options()
        file_paths, _ = QFileDialog.getOpenFileNames(None, "Select Files", "", "All Files (*)", options=options)
        if file_paths:
            self.fl_5.append("\n".join(file_paths))

    def get_6(self):
        options = QFileDialog.Options()
        file_paths, _ = QFileDialog.getOpenFileNames(None, "Select Files", "", "All Files (*)", options=options)
        if file_paths:
            self.fl_6.append("\n".join(file_paths))

    def get_outdir(self):
        options = QFileDialog.Options()
        directory_path = QFileDialog.getExistingDirectory(None, "Select Directory", options=options)
        if directory_path:
            self.outdir.setText(directory_path)


    def merge_files(self):
        file_names_text = self.fl_1.toPlainText()
        file_names_1 = file_names_text.split('\n')
        file_names_text = self.fl_2.toPlainText()
        file_names_2 = file_names_text.split('\n')
        file_names_text = self.fl_3.toPlainText()
        file_names_3 = file_names_text.split('\n')
        file_names_text = self.fl_4.toPlainText()
        file_names_4 = file_names_text.split('\n')
        file_names_text = self.fl_5.toPlainText()
        file_names_5 = file_names_text.split('\n')
        file_names_text = self.fl_6.toPlainText()
        file_names_6 = file_names_text.split('\n')

        text_line_1 =  self.led_1.text()
        text_line_2 =  self.led_1.text()
        text_line_3 =  self.led_1.text()
        text_line_4 =  self.led_1.text()
        text_line_5 =  self.led_1.text()
        text_line_6 =  self.led_1.text()

        output_directory = self.outdir.text()
        output_suffix = self.suffix.text()

        for fn_1, fn_2, fn_3, fn_4, fn_5, fn_6 in zip_longest(file_names_1, file_names_2, file_names_3, file_names_4, file_names_5, file_names_6, fillvalue=None):

            base_name = os.path.splitext(os.path.basename(fn_1))[0]
            new_file_name = f"{base_name}{output_suffix}.txt"
            output_file_name = os.path.join(output_directory, new_file_name)
            print("output_file_name: ",output_file_name)

            # Open for writing 
            with open(output_file_name, 'w', encoding='utf-8', errors='ignore') as Ofile:
                # Ofile.write("") All file writing occurs within this indentation loop
                if fn_1 is not None and os.path.exists(fn_1):
                    try:
                        with open(fn_1, 'r', encoding='utf-8', errors='ignore') as infile:
                            content = infile.read()
                            Ofile.write(content)
                            Ofile.write("\n\n")
                    except Exception as e:
                        print(f"An error occurred while reading the file: {e}")
                else:
                    print("The file is None or does not exist.")

                if text_line_1:
                    Ofile.write("\n")
                    Ofile.write(text_line_1)
                    Ofile.write("\n")

                if fn_2 is not None and os.path.exists(fn_2):
                    try:
                        with open(fn_2, 'r', encoding='utf-8', errors='ignore') as infile:
                            content = infile.read()
                            Ofile.write(content)
                            Ofile.write("\n\n")
                    except Exception as e:
                        print(f"An error occurred while reading the file: {e}")
                else:
                    print("The file is None or does not exist.")

                if text_line_2:
                    Ofile.write("\n")
                    Ofile.write(text_line_2)
                    Ofile.write("\n")

                if fn_3 is not None and os.path.exists(fn_3):
                    try:
                        with open(fn_3, 'r', encoding='utf-8', errors='ignore') as infile:
                            content = infile.read()
                            Ofile.write(content)
                            Ofile.write("\n\n")
                    except Exception as e:
                        print(f"An error occurred while reading the file: {e}")
                else:
                    print("The file is None or does not exist.")

                if text_line_3:
                    Ofile.write("\n")
                    Ofile.write(text_line_3)
                    Ofile.write("\n")

                if fn_4 is not None and os.path.exists(fn_4):
                    try:
                        with open(fn_4, 'r', encoding='utf-8', errors='ignore') as infile:
                            content = infile.read()
                            Ofile.write(content)
                            Ofile.write("\n\n")
                    except Exception as e:
                        print(f"An error occurred while reading the file: {e}")
                else:
                    print("The file is None or does not exist.")

                if text_line_4:
                    Ofile.write("\n")
                    Ofile.write(text_line_4)
                    Ofile.write("\n")

                if fn_5 is not None and os.path.exists(fn_5):
                    try:
                        with open(fn_5, 'r', encoding='utf-8', errors='ignore') as infile:
                            content = infile.read()
                            Ofile.write(content)
                            Ofile.write("\n\n")
                    except Exception as e:
                        print(f"An error occurred while reading the file: {e}")
                else:
                    print("The file is None or does not exist.")

                if text_line_5:
                    Ofile.write("\n")
                    Ofile.write(text_line_5)
                    Ofile.write("\n")

                if fn_6 is not None and os.path.exists(fn_6):
                    try:
                        with open(fn_6, 'r', encoding='utf-8', errors='ignore') as infile:
                            content = infile.read()
                            Ofile.write(content)
                            Ofile.write("\n\n")
                    except Exception as e:
                            print(f"An error occurred while reading the file: {e}")
                else:
                    print("The file is None or does not exist.")

                if text_line_6:
                    Ofile.write("\n")
                    Ofile.write(text_line_6)
                    Ofile.write("\n")

    def well_ok_then(self):
        # Create the QMessageBox
        msgBox = QMessageBox()
        msgBox.setIcon(QMessageBox.Warning)
        msgBox.setText("You are about to lose all unsaved changes.")
        msgBox.setInformativeText("This action cannot be undone. Do you wish to continue?")
        msgBox.setWindowTitle("Warning")
        msgBox.setStandardButtons(QMessageBox.Ok | QMessageBox.Cancel)

        # Execute the QMessageBox and check the response
        returnValue = msgBox.exec()
        if returnValue == QMessageBox.Ok:
            #print("OK clicked, proceeding with the action.")
            # Close the dialog or perform other clean up actions
            self.dialog.close()
        else:
            print("Cancel clicked, action aborted.")
            return

        


## ---------------------- Consolidate 




## --------------------- Editor Actions
class EDitorActions:
    def __init__(self, edit_1=None, edit_2=None, edit_3=None ):
        self.edit_1 = edit_1
        self.edit_2 = edit_2
        self.edit_3 = edit_3
        self.networkManager = None

#   OPEN EDIT-1
#
    def open_edit1(self):
        if self.edit_1:
            file_dialog = QFileDialog()
            file_path, _ = file_dialog.getOpenFileName(None, 'Open File', '', 'All Files (*)')
            if file_path:
                with open(file_path, 'r', encoding='utf-8', errors='ignore') as file:
                    self.edit_1.setPlainText(file.read())
        else:
            print("Edit_1 is not set.")

#   APPEND EDIT-1
#
    def append_multiple_1(self):
        # Store the current contents of edit_1
        original_content = self.edit_1.toPlainText()
        # Add a blank line to separate the original editor contents
        original_content += "\n"
        # Get the file names selected by the user
        file_names, _ = QFileDialog.getOpenFileNames(None, "Select Files", "", "All Files (*)")

        # Check if any file was selected
        if file_names:
            self.edit_1.clear()
            # Open each selected file and append its contents to the string
            for file_name in file_names:
                with open(file_name, 'r', encoding='utf-8', errors='ignore') as file:
                    original_content += file.read().strip()  # Append contents of each file
                    original_content += "\n\n"  # Add a blank line between each file's content

            # Store the combined string into edit_1
            self.edit_1.setPlainText(original_content)


#   SAVE EDIT-1
#
    def save_edit1(self):
        if self.edit_1:
            file_dialog = QFileDialog()
            file_path, _ = file_dialog.getSaveFileName(None, 'Save File', '', 'Text Files (*.txt);;All Files (*)')
            if file_path:
                with open(file_path, 'w', encoding='utf-8') as file:
                    file.write(self.edit_1.toPlainText())
        else:
            print("Edit_1 is not set.")

#   CLEAR EDIT-1
#
    def clear_edit1(self):
        self.edit_1.clear()
        default_format = QTextCharFormat()
        self.edit_1.setCurrentCharFormat(default_format)

#   SCAN TEXT
#   Searches for a list of terms in editor 1 within a list of directores in editor 2
#   Each .txt file in each directory is processed and all the match counts are presented
#   In editor 3 in csv (commma delimited) format ready for import into a spreadsheet
#
    def scan_text_files(self):
        search_terms = self.edit_1.toPlainText().splitlines()
        directory_paths = self.edit_2.toPlainText().splitlines()
        results = ""
        for directory in directory_paths:
            if os.path.isdir(directory):
                results += f"\n{directory}\n\n"
                for filename in os.listdir(directory):
                    if filename.endswith('.txt'):
                        file_path = os.path.join(directory, filename)
                        with open(file_path, 'r', encoding='utf-8') as file:
                            content = file.read()
                        
                        term_counts = {term: 0 for term in search_terms}

                        for term in search_terms:
                            term_counts[term] = content.count(term)

                        for term, count in term_counts.items():
                            if count > 0:
                                results += f"{filename}, {term}, {count}\n"

        self.edit_3.setPlainText(results)

#   SELECT DIRECTORIES
#   This is a companion method for SCAN TEXT that allows the user to select directories containing text files
#   It then adds these to the edit_2 editor list of directories.
#
    def select_directories(self):
        selected_directories = []

        while True:
            directory = QFileDialog.getExistingDirectory(None, "Select Directory")

            if directory:
                selected_directories.append(directory)
                txt_files = [f for f in os.listdir(directory) if f.endswith('.txt')]

                if txt_files:
                    file_list = "\n".join(txt_files)
                    QMessageBox.information(None, "Files in Directory", f"Directory: {directory}\n\nFiles:\n{file_list}")
                else:
                    QMessageBox.warning(None, "No Text Files", f"No .txt files found in the directory: {directory}")

            continue_selection = QMessageBox.question(None, "Continue Selection",
                                                      "Do you want to select another directory?",
                                                      QMessageBox.Yes | QMessageBox.No)

            if continue_selection == QMessageBox.No:
                break

        current_directories = self.edit_2.toPlainText().splitlines()
        updated_directories = current_directories + selected_directories
        self.edit_2.setPlainText("\n".join(updated_directories))

# Companion to SCAN TEXT
# Saves the contents edit_3 (the scan text report) as a CSV file for impor to Excel
#
    def save_as_csv(self):
        options = QFileDialog.Options()
        # options |= QFileDialog.DontUseNativeDialog
        file_path, _ = QFileDialog.getSaveFileName(None, "Save CSV File", "", "CSV Files (*.csv);;All Files (*)", options=options)
        
        if file_path:
            if not file_path.endswith('.csv'):
                file_path += '.csv'
            try:
                with open(file_path, 'w', encoding='utf-8', newline='') as file:
                    file.write(self.edit_3.toPlainText())
                QMessageBox.information(None, "Success", f"File saved successfully as {file_path}")
            except Exception as e:
                QMessageBox.critical(None, "Error", f"Failed to save file: {str(e)}")


# PDF Crusher
#
    def pdf_crusher(self):
        source_directory = QFileDialog.getExistingDirectory(None, 'Select Source Directory')
        output_directory = QFileDialog.getExistingDirectory(None, 'Select Output Directory')

        if source_directory and output_directory:
            for filename in os.listdir(source_directory):
                if filename.endswith('.pdf'):
                    file_path = os.path.join(source_directory, filename)
                    extracted_text = self.extract_text_from_pdf(file_path)
                    output_file_path = os.path.join(output_directory, os.path.splitext(filename)[0] + '.txt')
                    with open(output_file_path, 'w', encoding='utf-8', errors='ignore') as f:
                        f.write(extracted_text)
            QMessageBox.information(None, "Processing complete", "All PDF files have been processed successfully.")
        else:
            QMessageBox.information(None, "Operation cancelled", "No directory selected or operation cancelled.")

    def extract_text_from_pdf(self, file_path):
        text = []
        with fitz.open(file_path) as doc:
            for page in doc:
                text.append(page.get_text())

        return"\n".join(text)





#   OPEN EDIT-2
#
    def open_edit2(self):
        if self.edit_2:
            file_dialog = QFileDialog()
            file_path, _ = file_dialog.getOpenFileName(None, 'Open File', '', 'All Files (*)')
            if file_path:
                with open(file_path, 'r', encoding='utf-8', errors='ignore') as file:
                    self.edit_2.setPlainText(file.read())
        else:
            print("Edit_2 is not set.")


#   LOAD PPT to TEXT in EDIT-2
#

    def load_ppt_to_edit2(self):
        if self.edit_2:
            file_dialog = QFileDialog()
            # Update the filter to select only PowerPoint files
            file_path, _ = file_dialog.getOpenFileName(None, 'Open File', '', 'PowerPoint Files (*.ppt *.pptx)')
            if file_path:
                presentation = Presentation(file_path)
                all_text = []
                for slide_number, slide in enumerate(presentation.slides, start=1):
                    slide_text = [f"Slide: {slide_number}"]
                    # Extract text from slide shapes
                    for shape in slide.shapes:
                        if shape.has_text_frame:
                            for paragraph in shape.text_frame.paragraphs:
                                for run in paragraph.runs:
                                    slide_text.append(run.text)
                    # Extract text from speaker notes
                    if slide.notes_slide:
                        notes = slide.notes_slide.notes_text_frame
                        if notes:
                            notes_text = ['Speaker Notes:']
                            for paragraph in notes.paragraphs:
                                for run in paragraph.runs:
                                    notes_text.append(run.text)
                            slide_text.extend(notes_text)
                    all_text.extend(slide_text)
                # Join all text and set it to edit_2
                all_text_str = '\n'.join(all_text)
                self.edit_2.setPlainText(all_text_str)
        else:
            print("Edit_2 is not set.")


#   CONVERT all PPT files in a directory to TEXT 
#
    '''
    def Xppt_crusher(self):
        source_directory = QFileDialog.getExistingDirectory(None, 'Select Source Directory')
        output_directory = QFileDialog.getExistingDirectory(None, 'Select Output Directory')

        if source_directory and output_directory:
            for filename in os.listdir(source_directory):
                if filename.endswith('.pptx') and not filename.startswith('~$'):
                    file_path = os.path.join(source_directory, filename)
                    extracted_text = self.extract_text_from_pptx(file_path)
                    extracted_text = extracted_text.encode('utf-8', errors='ignore').decode('utf-8')
                    output_file_path = os.path.join(output_directory, os.path.splitext(filename)[0] + '.txt')
                    with open(output_file_path, 'w') as f:
                        f.write(extracted_text)
            print("Processing complete.")
        else:
            print("Operation cancelled or no directory selected.")

    def extract_text_from_pptx(self, file_path):
        prs = Presentation(file_path)
        text = []

        for slide in prs.slides:
            # Extract text from each shape in the slide
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text.append(shape.text)

            # Extract text from speaker notes
            if slide.has_notes_slide:
                notes_slide = slide.notes_slide
                if notes_slide.notes_text_frame:
                    for paragraph in notes_slide.notes_text_frame.paragraphs:
                        text.append(paragraph.text)

        return "\n".join(text)
    '''

    def ppt_crusher(self):
        source_directory = QFileDialog.getExistingDirectory(None, 'Select Source Directory')
        output_directory = QFileDialog.getExistingDirectory(None, 'Select Output Directory')

        if source_directory and output_directory:
            for filename in os.listdir(source_directory):
                if filename.endswith('.pptx'):
                    file_path = os.path.join(source_directory, filename)
                    extracted_text = self.extract_text_from_pptx(file_path)
                    output_file_path = os.path.join(output_directory, os.path.splitext(filename)[0] + '.txt')
                    with open(output_file_path, 'w', encoding='utf-8',errors='ignore') as f:
                        f.write(extracted_text)
            print("Processing complete.")
        else:
            print("Operation cancelled or no directory selected.")

    def extract_text_from_pptx(self, file_path):
        prs = Presentation(file_path)
        text = []

        for slide_number, slide in enumerate(prs.slides, start=1):
            slide_text = [f"Slide: {slide_number}"]
            # Extract text from each shape in the slide
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    slide_text.append(shape.text)
            # Extract text from speaker notes
            if slide.has_notes_slide:
                notes_slide = slide.notes_slide
                if notes_slide.notes_text_frame:
                    notes_text = ['Speaker Notes:']
                    for paragraph in notes_slide.notes_text_frame.paragraphs:
                        notes_text.append(paragraph.text)
                    slide_text.extend(notes_text)
            text.extend(slide_text)

        return "\n".join(text)
    


#   APPEND EDIT-2
#
    def append_multiple_2(self):
        # Store the current contents of edit_2
        original_content = self.edit_2.toPlainText()
        # Add a blank line to separate the original editor contents
        original_content += "\n"
        # Get the file names selected by the user
        file_names, _ = QFileDialog.getOpenFileNames(None, "Select Files", "", "All Files (*)")

        # Check if any file was selected
        if file_names:
            self.edit_2.clear()
            # Open each selected file and append its contents to the string
            for file_name in file_names:
                with open(file_name, 'r', encoding='utf-8', errors='ignore') as file:
                    original_content += file.read().strip()  # Append contents of each file
                    original_content += "\n\n"  # Add a blank line between each file's content

            # Store the combined string into edit_2
            self.edit_2.setPlainText(original_content)

#   SAVE EDIT-2
#
    def save_edit2(self):
        if self.edit_2:
            file_dialog = QFileDialog()
            file_path, _ = file_dialog.getSaveFileName(None, 'Save File', '', 'Text Files (*.txt);;All Files (*)')
            if file_path:
                with open(file_path, 'w', encoding='utf-8') as file:
                    file.write(self.edit_2.toPlainText())
        else:
            print("Edit_2 is not set.")


#   CLEAR EDIT-2
#
    def clear_edit2(self):
        self.edit_2.setPlainText('') ## Remove formatting from Markdown feature
        self.edit_2.clear()
        default_format = QTextCharFormat()
        self.edit_2.setCurrentCharFormat(default_format)

#   OPEN EDIT-3
#
    def open_edit3(self):
        if self.edit_3:
            file_dialog = QFileDialog()
            file_path, _ = file_dialog.getOpenFileName(None, 'Open File', '', 'All Files (*)')
            if file_path:
                with open(file_path, 'r', encoding='utf-8', errors='ignore') as file:
                    self.edit_3.setPlainText(file.read())
        else:
            print("Edit_3 is not set.")

#   APPEND EDIT-3
#
    def append_multiple_3(self):
        # Store the current contents of edit_2
        original_content = self.edit_3.toPlainText()
        # Add a blank line to separate the original editor contents
        original_content += "\n"
        # Get the file names selected by the user
        file_names, _ = QFileDialog.getOpenFileNames(None, "Select Files", "", "All Files (*)")

        # Check if any file was selected
        if file_names:
            self.edit_3.clear()
            # Open each selected file and append its contents to the string
            for file_name in file_names:
                with open(file_name, 'r', encoding='utf-8', errors='ignore') as file:
                    original_content += file.read().strip()  # Append contents of each file
                    original_content += "\n\n"  # Add a blank line between each file's content

            # Store the combined string into edit_2
            self.edit_3.setPlainText(original_content)

#   SAVE EDIT-3
#

    def save_edit3(self):
        if self.edit_3:
            file_dialog = QFileDialog()
            file_path, _ = file_dialog.getSaveFileName(None, 'Save File', '', 'Text Files (*.txt);;All Files (*)')
            if file_path:
                with open(file_path, 'w', encoding='utf-8') as file:
                    file.write(self.edit_3.toPlainText())
        else:
            print("Edit_3 is not set.")


#   CLEAR EDIT-3
#
    def clear_edit3(self):
        self.edit_3.clear()
        default_format = QTextCharFormat()
        self.edit_3.setCurrentCharFormat(default_format)


## --- Multiple editor methods

#   CLEAR EDIT-ALL
#
    def clear_editALL(self):
        self.edit_1.clear()
        self.edit_2.clear()
        self.edit_3.clear()

#   COPY 1 -> 2
#
    def copy_1_2(self):
        text_1 = self.edit_1.toPlainText()
        self.edit_2.setPlainText(text_1)

#   COPY 2 -> 1
#
    def copy_2_1(self):
        text_2 = self.edit_2.toPlainText()
        self.edit_1.setPlainText(text_2)

#   COPY 2 -> 3
#
    def copy_2_3(self):
        text_1 = self.edit_2.toPlainText()
        self.edit_3.setPlainText(text_1)

#   COPY 3 -> 2
#
    def copy_3_2(self):
        text_2 = self.edit_3.toPlainText()
        self.edit_2.setPlainText(text_2)



#   CHAT MODE >chat<
#
    def chat_step(self):
        text_1 = self.edit_1.toPlainText()
        text_2 = self.edit_2.toPlainText()
        text_3 = self.edit_3.toPlainText()

        # Clear all three text editors
        self.edit_1.clear()
        self.edit_2.clear()
        self.edit_3.clear()

        # Create history string
        history_string = f"{text_1}\nHuman: {text_2}\n\nAI: {text_3}\n"

        # Store the history string into edit_1
        self.edit_1.setPlainText(history_string)


#   Display Markdown code contained in EDIT_3 formatted in EDIT_2
#

    def display_as_markdown(self):
        # Get markdown content from edit_3
        markdown_content = self.edit_3.toPlainText()
        
        # Convert markdown to HTML
        html_content = markdown.markdown(markdown_content)
        
        # Display formatted markdown in edit_2
        self.edit_2.setHtml(html_content)

    def clipboard_edit1(self):
        self.edit_1.selectAll()
        self.edit_1.copy()
        cursor = self.edit_1.textCursor()
        cursor.clearSelection()
        self.edit_1.setTextCursor(cursor)

    def clipboard_edit2(self):
        self.edit_2.selectAll()
        self.edit_2.copy()
        cursor = self.edit_2.textCursor()
        cursor.clearSelection()
        self.edit_2.setTextCursor(cursor)

    def clipboard_edit3(self):
        self.edit_3.selectAll()
        self.edit_3.copy()
        cursor = self.edit_3.textCursor()
        cursor.clearSelection()
        self.edit_3.setTextCursor(cursor)

    def save_to_ppt(self, startup_location):
        # Use the provided template from the startup location
        template_file_path = os.path.join(startup_location, "ascend_template.pptx")
        
        if template_file_path:
            # Prompt the user for the output file name and location
            output_options = QFileDialog.Options()
            output_file_name, _ = QFileDialog.getSaveFileName(None, "Save PowerPoint", "", "PowerPoint Files (*.pptx);;All Files (*)", options=output_options)
            
            if output_file_name:
                # Copy the template file to the new filename and location
                shutil.copy(template_file_path, output_file_name)
                
                # Open the copied PowerPoint file for editing
                prs = Presentation(output_file_name)
                
                text = self.edit_3.toPlainText()
                
                # Adjust the split to capture titles and content including speaker notes
                slides_data = re.split(r'(Slide \d{1,2}: )', text)[1:]
                
                # Pair up the split results (slide marker with title, and content including speaker notes
                slides_text = [slides_data[n] + slides_data[n+1] for n in range(0, len(slides_data), 2)]

                for slide_text in slides_text:

                    # Splitting the slide text into lines
                    lines = slide_text.strip().split('\n')
    
                    # The first line is the title
                    title = lines[0].strip()
    
                    # Initialize variables to hold bullet points and speaker notes
                    bullet_points = []
                    speaker_notes = ""
    
                    # Flag to indicate that we've reached the speaker notes section
                    in_speaker_notes = False
    
                    # Iterate through the rest of the lines to separate bullet points and speaker notes
                    for line in lines[1:]:  # Skip the title
                        if line.startswith("SPEAKERNOTES:"):
                            in_speaker_notes = True
                            speaker_notes += line.replace("SPEAKERNOTES:", "").strip() + "\n"
                            continue
        
                        if in_speaker_notes:
                            speaker_notes += line.strip() + "\n"
                        else:
                            bullet_points.append(line.strip())
    
                    # Remove trailing newline from speaker notes if it exists
                    speaker_notes = speaker_notes.rstrip()                    

                    # Add a slide and determine the layout based on title content
                    if "TITLE: " in title:
                        slide_layout = prs.slide_layouts[1]  # Use layout 1 for slides with "TITLE: "
                        title = title.replace("TITLE: ", "")  # Remove "TITLE: " from the title
                        title = re.sub(r"(SLIDE \d{1,2}: |Slide \d{1,2}: )", "", title)
                        subtitle_index = 2  # Subtitle placeholder index in layout 1
                        slide = prs.slides.add_slide(slide_layout)
                        title_placeholder = slide.placeholders[1]
                        title_placeholder.text = title  # Set the title                       
                        subtitle_placeholder = slide.placeholders[subtitle_index]
                        # Combine bullet points into a single string for the subtitle
                        subtitle_text = "\n".join([bullet.lstrip('- ').strip() for bullet in bullet_points])
                        subtitle_placeholder.text = subtitle_text
                    elif "SECTION:" in title:
                        slide_layout = prs.slide_layouts[3]  # Use layout 3 for slides with "SECTION:"
                        title = title.replace("SECTION:", "")  # Remove "SECTION:" from the title
                        title = re.sub(r"(SLIDE \d{1,2}: |Slide \d{1,2}: )", "", title)
                        subtitle_index = 2  # Subtitle placeholder index in layout 3
                        slide = prs.slides.add_slide(slide_layout)
                        title_placeholder = slide.placeholders[1]
                        title_placeholder.text = title  # Set the title                       
                        subtitle_placeholder = slide.placeholders[subtitle_index]
                        # Combine bullet points into a single string for the subtitle
                        subtitle_text = "\n".join([bullet.lstrip('- ').strip() for bullet in bullet_points])
                        subtitle_placeholder.text = subtitle_text
                    else:
                        slide_layout = prs.slide_layouts[10]  # Default to layout 10 for 'Title and Content'
                        content_index = 2  # Content placeholder index in layout 10
                        slide = prs.slides.add_slide(slide_layout)                   
                        # For layouts with a content placeholder, add bullet points
                        title_placeholder = slide.placeholders[0]
                        title = re.sub(r"(SLIDE \d{1,2}: |Slide \d{1,2}: )", "", title)
                        title_placeholder.text = title  # Set the title
                        content_placeholder = slide.placeholders[content_index]
                        if content_placeholder.has_text_frame:
                            for bullet in bullet_points:
                                cleaned_bullet = bullet.lstrip('- ').strip()
                                if cleaned_bullet:
                                    p = content_placeholder.text_frame.add_paragraph()
                                    p.text = cleaned_bullet
                    
                    # Add speaker notes if they exist
                    if speaker_notes:
                        notes_slide = slide.notes_slide
                        notes_slide.notes_text_frame.text = speaker_notes
                
                # Save the presentation
                prs.save(output_file_name)                


    def enumerate_powerpoint_layouts(self): ## List Layouts
        options = QFileDialog.Options()
        # Set the dialog to open in the directory where the script is running, or specify a path
        initialDir = ""
        pptxFileName, _ = QFileDialog.getOpenFileName(None, "Open PowerPoint File", initialDir, "PowerPoint Files (*.pptx);;All Files (*)", options=options)
        
        if pptxFileName:
            prs = Presentation(pptxFileName)
            print(f"Layout information for: {pptxFileName}")
            for index, layout in enumerate(prs.slide_layouts):
                print(f"\nLayout {index}: {layout.name}")
                for placeholder in layout.placeholders:
                    print(f"  Placeholder {placeholder.placeholder_format.idx}: {placeholder.name}")
            # The presentation is automatically closed when the method exits and `prs` is deleted/garbage collected
        else:
            print("No file selected.")
    
    def fetch_url(self,max_web_chars=10000):
        self.max_web_chars = max_web_chars
        # Prompt user for URL
        url, ok = QInputDialog.getText(None, 'Input URL', 'Enter the URL:')
        
        if ok and url:
            # Create a QNetworkAccessManager and make a request
            self.networkManager = QNetworkAccessManager()
            self.networkManager.finished.connect(self._on_download_finished)
            self.networkManager.get(QNetworkRequest(QUrl(url)))

    def _on_download_finished(self, reply):
        # Check for errors
        if reply.error() == QNetworkReply.NoError:
            # Read response data
            content = reply.readAll().data().decode('utf-8')
            # Set content to QTextEdit ###
            soup = BeautifulSoup(content, "html.parser")
            text = soup.get_text()
            text = '\n'.join(line.strip() for line in text.splitlines() if line.strip())
            text = text[:self.max_web_chars]   ## Truncate the number of characters
            self.edit_2.setText(text)
        else:
            QMessageBox.warning(None, 'Error', 'Failed to download content: ' + reply.errorString())
        
        # Cleanup: Remove the QNetworkAccessManager instance 
        self.networkManager.deleteLater()



## ----------- Convert a directory of text files to PowerPoint slides.
        
    def hatch_ppt(self, startup_location):
        # Use the provided template from the startup location
        template_file_path = os.path.join(startup_location, "ascend_template.pptx")
        
        # Select directory instead of a single file
        dir_options = QFileDialog.Options()
        selected_dir = QFileDialog.getExistingDirectory(None, "Select Directory", options=dir_options)
        
        if selected_dir and template_file_path:
            for file in os.listdir(selected_dir):
                if file.endswith(".txt"):
                    text_file_path = os.path.join(selected_dir, file)
                    output_file_name = os.path.splitext(text_file_path)[0] + ".pptx"
                    
                    # Copy the template to the new PowerPoint file name
                    shutil.copy(template_file_path, output_file_name)
                    
                    # Open the new PowerPoint file for editing
                    prs = Presentation(output_file_name)
                    
                    # Read text file content
                    with open(text_file_path, 'r', encoding='utf-8') as f:
                        text = f.read()
                    
                    # Process text 
                    slides_data = re.split(r'(Slide \d{1,2}: )', text)[1:]
                    slides_text = [slides_data[n] + slides_data[n+1] for n in range(0, len(slides_data), 2)]

                    for slide_text in slides_text:
                        # Splitting the slide text into lines
                        lines = slide_text.strip().split('\n')
    
                        # The first line is the title
                        title = lines[0].strip()
    
                        # Initialize variables to hold bullet points and speaker notes
                        bullet_points = []
                        speaker_notes = ""
    
                        # Flag to indicate that we've reached the speaker notes section
                        in_speaker_notes = False
    
                        # Iterate through the rest of the lines to separate bullet points and speaker notes
                        for line in lines[1:]:  # Skip the title
                            if line.startswith("SPEAKERNOTES:"):
                                in_speaker_notes = True
                                speaker_notes += line.replace("SPEAKERNOTES:", "").strip() + "\n"
                                continue
        
                            if in_speaker_notes:
                                speaker_notes += line.strip() + "\n"
                            else:
                                bullet_points.append(line.strip())
    
                        # Remove trailing newline from speaker notes if it exists
                        speaker_notes = speaker_notes.rstrip()                    

                        # Add a slide and determine the layout based on title content
                        if "TITLE: " in title:
                            slide_layout = prs.slide_layouts[1]  # Use layout 1 for slides with "TITLE: "
                            title = title.replace("TITLE: ", "")  # Remove "TITLE: " from the title
                            title = re.sub(r"(SLIDE \d{1,2}: |Slide \d{1,2}: )", "", title)
                            subtitle_index = 2  # Subtitle placeholder index in layout 1
                            slide = prs.slides.add_slide(slide_layout)
                            title_placeholder = slide.placeholders[1]
                            title_placeholder.text = title  # Set the title                       
                            subtitle_placeholder = slide.placeholders[subtitle_index]
                            # Combine bullet points into a single string for the subtitle
                            subtitle_text = "\n".join([bullet.lstrip('- ').strip() for bullet in bullet_points])
                            subtitle_placeholder.text = subtitle_text
                        elif "SECTION:" in title:
                            slide_layout = prs.slide_layouts[3]  # Use layout 3 for slides with "SECTION:"
                            title = title.replace("SECTION:", "")  # Remove "SECTION:" from the title
                            title = re.sub(r"(SLIDE \d{1,2}: |Slide \d{1,2}: )", "", title)
                            subtitle_index = 2  # Subtitle placeholder index in layout 3
                            slide = prs.slides.add_slide(slide_layout)
                            title_placeholder = slide.placeholders[1]
                            title_placeholder.text = title  # Set the title                       
                            subtitle_placeholder = slide.placeholders[subtitle_index]
                            # Combine bullet points into a single string for the subtitle
                            subtitle_text = "\n".join([bullet.lstrip('- ').strip() for bullet in bullet_points])
                            subtitle_placeholder.text = subtitle_text
                        else:
                            slide_layout = prs.slide_layouts[10]  # Default to layout 10 for 'Title and Content'
                            content_index = 2  # Content placeholder index in layout 10
                            slide = prs.slides.add_slide(slide_layout)                   
                            # For layouts with a content placeholder, add bullet points
                            title_placeholder = slide.placeholders[0]
                            title = re.sub(r"(SLIDE \d{1,2}: |Slide \d{1,2}: )", "", title)
                            title_placeholder.text = title  # Set the title
                            content_placeholder = slide.placeholders[content_index]
                            if content_placeholder.has_text_frame:
                                for bullet in bullet_points:
                                    cleaned_bullet = bullet.lstrip('- ').strip()
                                    if cleaned_bullet:
                                        p = content_placeholder.text_frame.add_paragraph()
                                        p.text = cleaned_bullet
                    
                        # Add speaker notes if they exist
                        if speaker_notes:
                            notes_slide = slide.notes_slide
                            notes_slide.notes_text_frame.text = speaker_notes
                
                        # Save the presentation
                    prs.save(output_file_name)


    def fetch_file(self):
        # Open a file dialog to let the user select a file
        options = QFileDialog.Options()
        fileName, _ = QFileDialog.getOpenFileName(None, "Select a file", "", 
                                                  "All Files (*);;PDF Files (*.pdf);;Word Documents (*.docx);;Text Files (*.txt);;Excel Files (*.xlsx);;HTML Files (*.html);;CSV Files (*.csv);;Markdown Files (*.md);;JSON Files (*.json);;XML Files (*.xml);;RTF Files (*.rtf);;Video Text Tracks (*.vtt)",
                                                  options=options)
        if fileName:
            # Determine the file type based on the file extension

            if fileName.endswith('.pdf'):  ## PDF
                text = self.read_pdf(fileName)
                self.edit_2.setText(text)

            elif fileName.endswith('.docx'):  ## DOCX
                text = self.read_docx(fileName)
                self.edit_2.setText(text)
                # more code here

            elif fileName.endswith('.xlsx'): ## XLSX
                html = self.read_xlsx_to_html(fileName)
                self.edit_2.setHtml(html)  

            elif fileName.endswith('.txt'): ## TXT
                try:
                    with open(fileName, 'r', encoding='utf-8', errors='ignore') as file:
                        text = file.read()
                        self.edit_2.setText(text)
                except FileNotFoundError:
                    print("The TEXT file was not found.")
                except Exception as e:
                    print(f"An error occurred while reading the TEXT file: {e}")

            elif fileName.endswith('.html'): ## HTML
                html_content = ""
                try:
                    with open(fileName, 'r', encoding='utf-8', errors='ignore') as file:
                        html_content = file.read()
                        self.edit_2.setHtml(html_content)  # Use setHtml for HTML content
                except FileNotFoundError:
                    print("The HTML file was not found.")
                except Exception as e:
                    print(f"An error occurred while reading the HTML file: {e}")
                    
            elif fileName.endswith('.csv'): ## CSV
                csv_content = ""
                try:
                    with open(fileName, mode='r', encoding='utf-8', errors='ignore') as file:
                        reader = csv.reader(file)
                        for row in reader:
                            # Convert each row to a string and append to csv_content
                            csv_content += ', '.join(row) + '\n'
                except FileNotFoundError:
                    print("The CSV file was not found.")
                except Exception as e:
                    print(f"An error occurred while reading the CSV file: {e}")
                self.edit_2.setText(csv_content)

            elif fileName.endswith('.rtf'):
                rtf_content = ""
                try:
                    with open(fileName, 'r', encoding='utf-8', errors='ignore') as file:
                        rtf_content = file.read()
                        self.edit_2.setText(rtf_content)
                except FileNotFoundError:
                    print("The RTF file was not found.")
                except Exception as e:
                    print(f"An error occurred while reading the RTF file: {e}")

            elif fileName.endswith('.json'):
                try:
                    with open(fileName, 'r', encoding='utf-8', errors='ignore') as file:
                        content = json.load(file)
                        pretty_content = json.dumps(content, indent=4)
                        self.edit_2.setText(pretty_content)
                except FileNotFoundError:
                    print("The JSON file was not found.")
                except Exception as e:
                    print(f"An error occurred while reading the JSON file: {e}")

            elif fileName.endswith('.md'):
                print("Markdown file selected.")
                html_content = ""
                md_content = ""
                try:
                    with open(fileName, 'r', encoding='utf-8', errors='ignore') as file:
                        md_content = file.read()
                        html_content = markdown.markdown(md_content)
                except FileNotFoundError:
                    print("The Markdown file was not found.")
                except Exception as e:
                    print(f"An error occurred while reading the Markdown file: {e}")
                self.edit_2.setHtml(html_content)

            elif fileName.endswith('.xml'):
                xml_content = ""
                try:
                    with open(fileName, 'r', encoding='utf-8') as file:
                        xml_content = file.read()
                except FileNotFoundError:
                    print("The XML file was not found.")
                except Exception as e:
                    print(f"An error occurred while reading the XML file: {e}")
                self.edit_2.setText(xml_content)

            elif fileName.endswith('.vtt'): ## VTT
                try:
                    with open(fileName, 'r', encoding='utf-8', errors='ignore') as file:
                        text = file.read()
                        self.edit_2.setText(text)
                except FileNotFoundError:
                    print("The VTT file was not found.")
                except Exception as e:
                    print(f"An error occurred while reading the VTT file: {e}")

            else:
                print("File type not recognized.")
                # more code here

    def read_docx(self, docxPath):
        # Open the .docx file
        doc = Document(docxPath)
        fullText = []
        for para in doc.paragraphs:
            fullText.append(para.text)
        return '\n'.join(fullText)
    
    def read_pdf(self, pdfPath):
        # Open the PDF file
        doc = fitz.open(pdfPath)
        text = ''
        # Extract text from each page
        for page in doc:
            text += page.get_text()
        return text
    
    def read_xlsx_to_html(self, xlsxPath):
        # Load the workbook and select the first sheet
        wb = openpyxl.load_workbook(xlsxPath)
        sheet = wb[wb.sheetnames[0]]
        html = "<table border='1'>"
        
        # Convert each row in the sheet to an HTML table row
        for row in sheet.iter_rows(values_only=True):
            html += "<tr>" + "".join(f"<td>{value}</td>" if value is not None else "<td></td>" for value in row) + "</tr>"
        
        html += "</table>"
        return html

## ---- File Merge

    def file_merge(self):
        # for item in self.settings:
        #    print("DEBUG-BEFORE:",item)
        
        dialog = FileMergeDialog()
        if dialog.exec_():
            print("\nSettings updated.")
        
        #for item in self.settings:
        #    print("DEBUG-AFTER:",item)



    def file_lister(self):
        directory = QFileDialog.getExistingDirectory(None, "Select Directory")
        if directory:
            filenames = []
            tree = []
            for root, dirs, files in os.walk(directory):
                level = root.replace(directory, '').count(os.sep)
                indent = ' ' * 4 * level
                # Use + for directories in the ASCII tree
                tree.append(f"{indent}+ {os.path.basename(root)}/")  # Adding graphical tree connections
                subindent = ' ' * 4 * (level + 1)

                # Format directory names for self.edit_2
                dir_line = f"DIRECTORY: {os.path.basename(root)}"
                if level == 0:
                    filenames.append(dir_line)  # Add the base directory at the top
                else:
                    filenames.extend(['', dir_line])  # Skip one line before each new directory

                for f in files:
                    # Use | for files in the ASCII tree
                    tree.append(f"{subindent}| {f}")
                    filenames.append(f)  # List files unindented under their directory

            self.edit_2.setText("\n".join(filenames))
            self.edit_3.setText("\n".join(tree))
        else:
            self.edit_2.setText("No directory selected.")
            self.edit_3.setText("No directory selected.")



